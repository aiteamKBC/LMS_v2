"""AI flash-card generation for the deck builder.

Kept in its own module (like learner_api splits helpers out of views) so the
views stay thin. This is the engagement app's own copy of the text-extraction +
OpenAI-call pieces — it deliberately does NOT import quiz_api, which is outside
this feature's scope.

Flow (mirrors quiz_api.generate_ai_questions, but flash-card shaped):
  1. extract_text_from_files() pulls plain text out of uploaded source files.
  2. generate_flashcards() builds a prompt (embedded core rules + the author's
     editable instructions from the UI) and asks OpenAI for a strict-JSON list
     of {question, answer, category, difficulty} cards.

Points/targeting live elsewhere; this module only produces card content.
"""
import re
from pathlib import Path

from django.conf import settings

# Upload / content guardrails (same order of magnitude as the quiz builder).
MAX_FILE_SIZE = 50 * 1024 * 1024          # 50 MB per file
MAX_CONTENT_CHARS = 60000                 # total source text sent to the model
MAX_FILES = 10

DEFAULT_CARD_COUNT = 25
MIN_CARD_COUNT = 1
MAX_CARD_COUNT = 60

DIFFICULTIES = ('easy', 'medium', 'hard')
TEXT_EXTENSIONS = {'.txt', '.md', '.csv', '.xml', '.html', '.htm', '.json'}

# The core prompt is embedded here in the backend. The builder UI shows an
# editable copy of the "author instructions" and sends them as
# custom_instructions, which are appended below — but they can never override
# the JSON-schema / source-grounding / safety rules.
EMBEDDED_FLASHCARD_PROMPT = """\
Act like an assessment-design SME in adult learning building high-quality revision \
flash cards for a UK workplace apprenticeship LMS.

You are creating a deck of flash cards a learner reviews AFTER a lecture to reinforce \
the material and earn engagement points. Each card is a single question/answer pair:
- "question" is the FRONT: a clear, self-contained prompt (a question, term, or scenario).
- "answer" is the BACK: a concise, correct answer — ideally 1-3 sentences. No preamble like \
"The answer is".

Rules:
- Generate exactly the requested number of cards when the source content supports it.
- Use ONLY the supplied source content / topic. Do not invent facts or use outside knowledge \
beyond what is needed to explain the source material.
- Cover distinct concepts across the material; do not write several cards testing the same idea.
- Progress difficulty across the deck (easy recall -> medium understanding -> harder application). \
Label each card's "difficulty" as easy, medium or hard.
- Give each card a short "category" (a topic/theme drawn from the material, e.g. "Scheduling", \
"Risk", "Branding").
- Where the material implies KSBs (Knowledge, Skills, Behaviours), you may reference them briefly \
in the answer, but keep answers focused and readable.
- Write for adult UK workplace learners; keep language clear and professional.
- Do not mention "the source", "the slides", "the document", or "the PDF" in the card text.
- Return structured JSON only, matching the provided schema."""


def _decode_text(raw):
    for encoding in ('utf-8-sig', 'utf-8', 'cp1252', 'latin-1'):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode('utf-8-sig', errors='replace')


def _extract_one(uploaded_file):
    """Best-effort plain-text extraction from a single uploaded file.

    Returns '' for files we can't read (rather than raising) so one bad file
    doesn't sink the whole batch; raises ValueError only for hard limits
    (oversized / unsupported type) that the caller should surface.
    """
    if uploaded_file.size > MAX_FILE_SIZE:
        raise ValueError(f'{uploaded_file.name} is larger than 50 MB.')

    extension = Path(uploaded_file.name).suffix.lower()
    uploaded_file.seek(0)

    if extension in TEXT_EXTENSIONS:
        return _decode_text(uploaded_file.read())

    if extension == '.pdf':
        from pypdf import PdfReader

        reader = PdfReader(uploaded_file)
        lines = []
        for index, page in enumerate(reader.pages[:40], start=1):
            text = re.sub(r'\s+', ' ', page.extract_text() or '').strip()
            if text:
                lines.append(f'Page {index}: {text}')
        return '\n'.join(lines)

    if extension == '.docx':
        from docx import Document

        document = Document(uploaded_file)
        lines = [re.sub(r'\s+', ' ', p.text).strip() for p in document.paragraphs if p.text.strip()]
        for table in document.tables:
            for row in table.rows:
                cells = [re.sub(r'\s+', ' ', c.text).strip() for c in row.cells if c.text.strip()]
                if cells:
                    lines.append(' | '.join(cells))
        return '\n'.join(lines)

    if extension in {'.pptx', '.pptm'}:
        from pptx import Presentation

        presentation = Presentation(uploaded_file)
        lines = []
        for index, slide in enumerate(presentation.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    text = re.sub(r'\s+', ' ', shape.text).strip()
                    if text:
                        parts.append(text)
            if parts:
                lines.append(f'Slide {index}: ' + ' | '.join(parts))
        return '\n'.join(lines)

    if extension in {'.xlsx', '.xlsm'}:
        from openpyxl import load_workbook

        workbook = load_workbook(uploaded_file, read_only=True, data_only=True)
        lines = []
        for sheet in workbook.worksheets[:3]:
            lines.append(f'Sheet: {sheet.title}')
            for row in sheet.iter_rows(values_only=True):
                values = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if values:
                    lines.append(' | '.join(values))
        return '\n'.join(lines)

    raise ValueError(
        f'Unsupported source file "{uploaded_file.name}". '
        'Upload TXT, MD, CSV, PDF, DOCX, PPTX or XLSX.'
    )


def extract_text_from_files(files):
    """Extract text from up to MAX_FILES uploads.

    Returns (source_text, readable_names, unreadable_names). Each file's text is
    prefixed with a "Source file:" header and the total is budget-capped so no
    single file dominates the model's context.
    """
    chunks = []
    readable = []
    unreadable = []
    for uploaded_file in list(files)[:MAX_FILES]:
        text = (_extract_one(uploaded_file) or '').strip()
        if text:
            chunks.append((uploaded_file.name, text))
            readable.append(uploaded_file.name)
        else:
            unreadable.append(uploaded_file.name)

    if not chunks:
        return '', readable, unreadable

    per_file_budget = max(4000, MAX_CONTENT_CHARS // len(chunks))
    joined = '\n\n'.join(f'Source file: {name}\n{text[:per_file_budget]}' for name, text in chunks)
    return joined[:MAX_CONTENT_CHARS], readable, unreadable


def clamp_count(value):
    try:
        count = int(value)
    except (TypeError, ValueError):
        count = DEFAULT_CARD_COUNT
    return max(MIN_CARD_COUNT, min(count, MAX_CARD_COUNT))


def _normalise_cards(data):
    raw = data.get('cards', []) if isinstance(data, dict) else []
    cards = []
    for item in raw:
        if not isinstance(item, dict):
            continue
        question = str(item.get('question') or item.get('front') or '').strip()
        answer = str(item.get('answer') or item.get('back') or '').strip()
        if not question or not answer:
            continue
        difficulty = str(item.get('difficulty') or 'medium').strip().lower()
        if difficulty not in DIFFICULTIES:
            difficulty = 'medium'
        cards.append({
            'question': question,
            'answer': answer,
            'category': str(item.get('category') or '').strip(),
            'difficulty': difficulty,
        })
    return cards


def generate_flashcards(source_text, topic='', custom_instructions='',
                        programme='', module='', count=DEFAULT_CARD_COUNT):
    """Call OpenAI and return a list of normalised card dicts.

    Raises RuntimeError with a user-facing message on any failure (missing key,
    package not installed, quota, invalid key, bad response) so the view can map
    it to an HTTP status.
    """
    if not settings.OPENAI_API_KEY:
        raise RuntimeError('OPENAI_API_KEY is not configured in backend .env')

    count = clamp_count(count)
    source_text = (source_text or '')[:MAX_CONTENT_CHARS]
    if not source_text and not topic:
        raise RuntimeError('No source content. Add a topic, paste lesson text, or upload readable files.')

    schema = {
        'type': 'object',
        'additionalProperties': False,
        'properties': {
            'cards': {
                'type': 'array',
                'minItems': 1,
                'maxItems': count,
                'items': {
                    'type': 'object',
                    'additionalProperties': False,
                    'properties': {
                        'question': {'type': 'string'},
                        'answer': {'type': 'string'},
                        'category': {'type': 'string'},
                        'difficulty': {'type': 'string', 'enum': list(DIFFICULTIES)},
                    },
                    'required': ['question', 'answer', 'category', 'difficulty'],
                },
            }
        },
        'required': ['cards'],
    }

    prompt = f"""{EMBEDDED_FLASHCARD_PROMPT}

Programme: {programme or 'Not provided'}
Module: {module or 'Not provided'}
Topic: {topic or 'Use the source content'}
Number of cards: {count}

Source content:
{source_text or topic}

Editable author instructions from the builder UI:
{custom_instructions or 'None provided. Follow the embedded core rules only.'}
"""

    try:
        from openai import OpenAI
    except ImportError:
        raise RuntimeError('OpenAI Python package is not installed. Run pip install -r requirments.txt.')

    try:
        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        response = client.responses.create(
            model=settings.OPENAI_MODEL,
            input=[
                {'role': 'system', 'content': 'You generate revision flash cards and return only structured JSON.'},
                {'role': 'user', 'content': prompt},
            ],
            text={
                'format': {
                    'type': 'json_schema',
                    'name': 'generated_flash_cards',
                    'schema': schema,
                    'strict': True,
                }
            },
        )
    except Exception as exc:  # noqa: BLE001 - map provider errors to messages
        message = str(exc)
        if 'insufficient_quota' in message or 'exceeded your current quota' in message or 'Error code: 429' in message:
            raise RuntimeError('OpenAI quota is exhausted for the configured API key. Check billing/quota or replace OPENAI_API_KEY.')
        if 'invalid_api_key' in message or 'Incorrect API key' in message:
            raise RuntimeError('The configured OpenAI API key is invalid. Update OPENAI_API_KEY in backend .env.')
        raise RuntimeError(f'Flash-card generation failed: {message}')

    import json

    try:
        data = json.loads(response.output_text)
    except (json.JSONDecodeError, TypeError):
        raise RuntimeError('AI returned invalid JSON.')

    cards = _normalise_cards(data)
    if not cards:
        raise RuntimeError('AI did not return usable cards.')
    return cards
