"""Rendering an uploaded deck into the slide model the learner page draws.

The decks here are built with python-pptx so the expected geometry and colours
are known exactly, which is the only way to assert on a renderer whose output is
otherwise "looks about right". Two behaviours get the most attention because
both produced visibly broken slides during development:

* the master's colour map, which a dark template inverts — ignoring it renders
  light body text as near-black on a near-black background;
* line spacing, which PowerPoint stores either as a multiple or as an exact
  length, and Length subclasses int.
"""
import json
from pathlib import Path
from tempfile import TemporaryDirectory

from django.test import Client, SimpleTestCase, override_settings

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from . import pptx_slides
from .views import COMPONENT_UPLOAD_ROOT


def build_deck(path, *, dark_master=False):
    """A one-slide 16:9 deck with a title, a filled box and a table."""
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # blank

    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1.5))
    paragraph = box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = 'Hello slides'
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x11, 0x22, 0x33)

    table = slide.shapes.add_table(2, 2, Inches(1), Inches(4), Inches(6), Inches(2)).table
    table.cell(0, 0).text = 'Left'
    table.cell(0, 1).text = 'Right'

    if dark_master:
        # What a dark template does: keep the light theme's slots but point the
        # text/background tokens at the opposite ones.
        color_map = presentation.slide_masters[0]._element.find(qn('p:clrMap'))
        color_map.set('tx1', 'lt1')
        color_map.set('bg1', 'dk1')
    presentation.save(str(path))
    return presentation


class SlideModelTests(SimpleTestCase):
    def setUp(self):
        self._directory = TemporaryDirectory()
        self.addCleanup(self._directory.cleanup)
        self.root = Path(self._directory.name)
        self.deck_path = self.root / 'deck.pptx'
        build_deck(self.deck_path)

    def render(self, path=None):
        return pptx_slides._deck_model(path or self.deck_path, self.root / 'assets', '/assets')

    def test_stage_matches_the_decks_own_dimensions(self):
        deck = self.render()
        self.assertEqual(deck['slideCount'], 1)
        # 13.333in x 7.5in at 96 dpi.
        self.assertEqual((deck['slideWidthPx'], deck['slideHeightPx']), (1280, 720))

    def test_text_carries_position_size_and_colour(self):
        shape = next(
            entry for entry in self.render()['slides'][0]['shapes']
            if entry['kind'] == 'text' and entry['paragraphs']
        )
        self.assertEqual((shape['x'], shape['y']), (96.0, 96.0))
        run = shape['paragraphs'][0]['runs'][0]
        self.assertEqual(run['text'], 'Hello slides')
        self.assertEqual(run['color'], '#112233')
        self.assertTrue(run['bold'])
        self.assertAlmostEqual(run['sizePx'], 48.0, places=1)  # 36pt

    def test_tables_keep_their_cells_and_column_shares(self):
        table = next(entry for entry in self.render()['slides'][0]['shapes'] if entry['kind'] == 'table')
        self.assertEqual(len(table['rows']), 2)
        self.assertEqual(
            [cell['paragraphs'][0]['runs'][0]['text'] for cell in table['rows'][0]],
            ['Left', 'Right'],
        )
        self.assertAlmostEqual(sum(table['colFractions']), 1.0, places=3)

    def test_dark_master_colour_map_gives_light_default_text(self):
        """The bug this guards: black text on a black slide.

        A dark template's clrMap points tx1 at the light slot. Resolving tx1
        against the default (light) mapping instead returns near-black, and every
        slide that leaves its body text to the master renders unreadable.
        """
        dark_path = self.root / 'dark.pptx'
        build_deck(dark_path, dark_master=True)
        light = self.render()['slides'][0]['shapes']
        dark = self.render(dark_path)['slides'][0]['shapes']

        def default_colors(shapes):
            return {entry['defaultTextColor'] for entry in shapes if entry.get('defaultTextColor')}

        self.assertEqual(default_colors(light), {'#000000'})
        self.assertEqual(default_colors(dark), {'#ffffff'})


class ColourAndSpacingTests(SimpleTestCase):
    def test_alpha_becomes_an_rgba_fill(self):
        self.assertEqual(pptx_slides._with_alpha('#ffffff', None), '#ffffff')
        self.assertEqual(pptx_slides._with_alpha('#ffffff', 0.15), 'rgba(255, 255, 255, 0.15)')

    def test_opaque_hex_rejects_translucent_colours(self):
        # Contrast has to be measured against something solid; a translucent
        # fill lets the slide behind it through instead.
        self.assertEqual(pptx_slides.opaque_hex('#102030'), '#102030')
        self.assertIsNone(pptx_slides.opaque_hex('rgba(0, 0, 0, 0.2)'))
        self.assertIsNone(pptx_slides.opaque_hex(None))

    def test_contrast_text_follows_the_surface(self):
        self.assertEqual(pptx_slides._contrast_text('#ffffff'), pptx_slides.DARK_TEXT)
        self.assertEqual(pptx_slides._contrast_text('#002615'), pptx_slides.LIGHT_TEXT)

    def test_svg_numbers_never_use_scientific_notation(self):
        # Path coordinates are EMU-scale; "%g" would emit 1.28e+07, which SVG
        # path data rejects, and the whole outline disappears.
        self.assertEqual(pptx_slides._svg_number(12808135), '12808135')
        self.assertEqual(pptx_slides._svg_number(1.256), '1.26')
        self.assertNotIn('e', pptx_slides._svg_number(1.28081e7))

    def test_brightness_lightens_and_darkens(self):
        self.assertEqual(pptx_slides._apply_brightness('#808080', 0.5), '#c0c0c0')
        self.assertEqual(pptx_slides._apply_brightness('#808080', -0.5), '#404040')
        self.assertEqual(pptx_slides._apply_brightness('#808080', 0), '#808080')


class RenderCacheTests(SimpleTestCase):
    def setUp(self):
        self._directory = TemporaryDirectory()
        self.addCleanup(self._directory.cleanup)
        self.media = Path(self._directory.name)
        self.relative = f'{COMPONENT_UPLOAD_ROOT}/MOD/COMP/deck.pptx'
        self.source = self.media / self.relative
        self.source.parent.mkdir(parents=True, exist_ok=True)
        build_deck(self.source)

    def render(self):
        with override_settings(MEDIA_ROOT=self.media):
            return pptx_slides.render_uploaded_deck(self.relative, self.source)

    def test_second_render_is_served_from_the_manifest(self):
        first = self.render()
        manifest = next((self.media / COMPONENT_UPLOAD_ROOT / pptx_slides.RENDER_DIR_NAME).glob('*/deck.json'))
        stored = json.loads(manifest.read_text(encoding='utf-8'))
        self.assertEqual(stored['source'], self.relative)
        manifest.write_text(json.dumps({**stored, 'deck': {'slideCount': 99}}), encoding='utf-8')
        # Same stamp, so the (doctored) cache is trusted rather than re-parsed.
        self.assertEqual(self.render()['slideCount'], 99)
        self.assertEqual(first['slideCount'], 1)

    def test_replacing_the_file_discards_the_old_render(self):
        self.render()
        build_deck(self.source)  # a new upload at the same path
        import os
        os.utime(self.source, (0, 0))
        self.assertEqual(self.render()['slideCount'], 1)

    def test_legacy_powerpoint_is_refused_with_an_explanation(self):
        legacy = self.media / 'deck.ppt'
        legacy.write_bytes(b'not really a deck')
        with self.assertRaises(pptx_slides.UnsupportedDeck) as caught:
            pptx_slides.render_uploaded_deck('deck.ppt', legacy)
        self.assertIn('97-2003', str(caught.exception))

    def test_an_unreadable_deck_reports_rather_than_raising_a_parse_error(self):
        broken = self.media / 'broken.pptx'
        broken.write_bytes(b'PK\x03\x04 truncated')
        with self.assertRaises(pptx_slides.UnsupportedDeck):
            pptx_slides.render_uploaded_deck('broken.pptx', broken)


class SlidesEndpointTests(SimpleTestCase):
    def setUp(self):
        self._directory = TemporaryDirectory()
        self.addCleanup(self._directory.cleanup)
        self.media = Path(self._directory.name)
        self.relative = f'{COMPONENT_UPLOAD_ROOT}/MOD/COMP/deck.pptx'
        source = self.media / self.relative
        source.parent.mkdir(parents=True, exist_ok=True)
        build_deck(source)
        self.client = Client()
        self.url = '/curriculum_api/curriculum/presentations/slides/'

    def get(self, src):
        with override_settings(MEDIA_ROOT=self.media):
            return self.client.get(self.url, {'src': src})

    def test_renders_the_deck_for_its_upload_path(self):
        response = self.get('/curriculum_api/curriculum/uploads/MOD/COMP/deck.pptx')
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()['slideCount'], 1)

    def test_accepts_an_absolute_url_for_the_same_upload(self):
        response = self.get('https://lms.example.test/curriculum_api/curriculum/uploads/MOD/COMP/deck.pptx')
        self.assertEqual(response.status_code, 200)

    def test_missing_file_explains_itself_instead_of_500ing(self):
        response = self.get('/curriculum_api/curriculum/uploads/MOD/COMP/gone.pptx')
        self.assertEqual(response.status_code, 404)
        self.assertIn('Re-upload', response.json()['error'])

    def test_a_path_outside_the_uploads_route_is_rejected(self):
        for src in ('', '/media/secrets.pptx', '/curriculum_api/curriculum/uploads/../../secrets.pptx'):
            self.assertEqual(self.get(src).status_code, 400, src)

    def test_rendered_assets_are_not_themselves_renderable(self):
        # The render cache is output; treating it as input would let a request
        # walk back into the renderer's own directory.
        src = f'/curriculum_api/curriculum/uploads/{pptx_slides.RENDER_DIR_NAME}/abc/deck.pptx'
        self.assertEqual(self.get(src).status_code, 400)


class EmuGeometryTests(SimpleTestCase):
    def test_group_children_are_mapped_into_the_groups_box(self):
        """A group scales its children from its own coordinate space."""
        class FakeGroup:
            left, top, width, height = Emu(1000), Emu(2000), Emu(4000), Emu(4000)

            class element:
                @staticmethod
                def find(_tag):
                    return None
        frame = pptx_slides.Frame().nested(FakeGroup())
        # With no chOff/chExt to read, children are simply offset by the group.
        self.assertEqual(frame.px(0, 0, 0, 0)['x'], round(1000 / pptx_slides.EMU_PER_PX, 2))
        self.assertEqual(frame.px(0, 0, 0, 0)['y'], round(2000 / pptx_slides.EMU_PER_PX, 2))


class PdfPageRenderTests(SimpleTestCase):
    """A PDF is rendered to page images rather than handed to the browser.

    Whether a browser previews a PDF in an iframe or offers it as a download is
    a setting on the reader's own machine — learners were getting the download
    prompt inside the frame, which no response header can override.
    """

    def setUp(self):
        self._directory = TemporaryDirectory()
        self.addCleanup(self._directory.cleanup)
        self.media = Path(self._directory.name)
        self.pdf = self.media / 'handout.pdf'
        self._write_pdf(self.pdf, pages=3)

    @staticmethod
    def _write_pdf(path, pages=1):
        import fitz
        document = fitz.open()
        for number in range(pages):
            page = document.new_page(width=612, height=792)  # US Letter
            page.insert_text((72, 144), f'Page {number + 1}', fontsize=24)
        document.save(str(path))
        document.close()

    def render(self, path=None):
        return pptx_slides._pdf_model(path or self.pdf, self.media / 'assets', '/assets')

    def test_each_page_becomes_one_full_page_image(self):
        model = self.render()

        self.assertEqual(model['slideCount'], 3)
        self.assertEqual(model['unit'], 'page', 'the viewer says "Page 1 of 3", not "Slide"')
        for index, page in enumerate(model['slides'], start=1):
            self.assertEqual(page['number'], index)
            self.assertEqual(len(page['shapes']), 1)
            shape = page['shapes'][0]
            self.assertEqual(shape['kind'], 'image')
            self.assertEqual((shape['x'], shape['y']), (0, 0))
            self.assertTrue(shape['src'].startswith('/assets/'))

    def test_the_stage_is_the_page_at_the_render_scale(self):
        model = self.render()

        # 612x792pt at 1.5x -> 918x1188 device pixels.
        self.assertEqual(model['slideWidthPx'], round(612 * pptx_slides.PDF_RENDER_ZOOM))
        self.assertEqual(model['slideHeightPx'], round(792 * pptx_slides.PDF_RENDER_ZOOM))

    def test_the_images_are_written_where_they_will_be_served_from(self):
        self.render()

        written = sorted(p.name for p in (self.media / 'assets').iterdir())
        self.assertEqual(len(written), 3, written)
        self.assertTrue(all(name.endswith('.png') for name in written), written)

    def test_a_long_document_is_capped_and_says_so(self):
        long_pdf = self.media / 'long.pdf'
        self._write_pdf(long_pdf, pages=pptx_slides.PDF_MAX_PAGES + 4)
        model = self.render(long_pdf)

        self.assertEqual(model['slideCount'], pptx_slides.PDF_MAX_PAGES)
        self.assertTrue(model['truncated'])
        self.assertEqual(model['totalPages'], pptx_slides.PDF_MAX_PAGES + 4)

    def test_a_short_document_is_not_marked_truncated(self):
        model = self.render()

        self.assertFalse(model['truncated'])
        self.assertEqual(model['totalPages'], 3)

    def test_a_pdf_goes_through_the_same_cache_as_a_deck(self):
        relative = f'{COMPONENT_UPLOAD_ROOT}/MOD/COMP/handout.pdf'
        source = self.media / relative
        source.parent.mkdir(parents=True, exist_ok=True)
        self._write_pdf(source, pages=2)
        with override_settings(MEDIA_ROOT=self.media):
            first = pptx_slides.render_uploaded_deck(relative, source)
            # Same stamp, so the second call must not re-rasterise.
            manifest = next(
                (self.media / COMPONENT_UPLOAD_ROOT / pptx_slides.RENDER_DIR_NAME).glob('*/deck.json')
            )
            stored = json.loads(manifest.read_text(encoding='utf-8'))
            manifest.write_text(json.dumps({**stored, 'deck': {'slideCount': 99}}), encoding='utf-8')
            second = pptx_slides.render_uploaded_deck(relative, source)

        self.assertEqual(first['slideCount'], 2)
        self.assertEqual(second['slideCount'], 99, 'the render was not read from the cache')

    def test_a_file_that_is_not_really_a_pdf_reports_rather_than_crashing(self):
        # A zero-page PDF cannot be written at all, so the realistic failure is
        # a truncated or mislabelled upload.
        relative = f'{COMPONENT_UPLOAD_ROOT}/MOD/COMP/broken.pdf'
        broken = self.media / relative
        broken.parent.mkdir(parents=True, exist_ok=True)
        broken.write_bytes(b'%PDF-1.7 truncated right here')

        with override_settings(MEDIA_ROOT=self.media):
            with self.assertRaises(pptx_slides.UnsupportedDeck) as caught:
                pptx_slides.render_uploaded_deck(relative, broken)
        self.assertIn('document', str(caught.exception).lower())

    def test_a_pdf_is_recognised_as_renderable(self):
        self.assertTrue(pptx_slides.is_renderable_deck_path('/uploads/a/b/handout.pdf'))
        self.assertTrue(pptx_slides.is_renderable_deck_path('/uploads/a/b/deck.pptx'))
        self.assertFalse(pptx_slides.is_renderable_deck_path('/uploads/a/b/notes.docx'))
