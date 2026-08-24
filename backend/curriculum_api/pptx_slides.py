"""Render an uploaded PowerPoint deck into a JSON slide model the browser can draw.

Why this exists: an uploaded deck has no inline preview of its own. Microsoft's
Office Online viewer can only show a file it can download from the public
internet, which rules out every private or local deployment, and hands the file
to a third party besides. There is no pure-Python PowerPoint rasteriser and no
LibreOffice on the box, so this module reads the OOXML with python-pptx and
emits geometry + text + images — a "good enough to read the slides" model that
the SlideDeckViewer positions absolutely on a fixed-size stage.

What is modelled: slide/layout/master backgrounds, the decorative shapes a
template puts on its master and layouts, text frames (runs with size, weight,
colour, alignment, bullets), pictures, tables, and grouped shapes. What is not:
charts, SmartArt and embedded objects become labelled placeholder boxes; effects
(shadows, 3-D), transitions and animations are dropped. Gradients collapse to
their first stop.

Renders are cached on disk beside the upload, keyed by the source file's size
and mtime, so a deck is parsed once and re-parsed after it is replaced.
"""

from __future__ import annotations

import hashlib
import json
import logging
import math
import shutil
from pathlib import Path

from django.conf import settings

from lxml import etree
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.text.text import _Run
from pptx.util import Length

logger = logging.getLogger(__name__)

# Bumped when the emitted model changes shape, so cached renders from an older
# version of this module are discarded instead of being replayed.
MODEL_VERSION = 6

#: Cache directory name, a sibling of the uploads it renders.
RENDER_DIR_NAME = '_slide_renders'

#: python-pptx reads OOXML only; legacy binary PowerPoint is a different format.
OOXML_SUFFIXES = {'.pptx', '.ppsx', '.pptm', '.ppsm'}
LEGACY_SUFFIXES = {'.ppt', '.pps'}

EMU_PER_PX = 9525          # 914400 EMU/inch / 96 px/inch
PT_TO_PX = 96 / 72

DEFAULT_SLIDE_WIDTH_PX = 960
DEFAULT_SLIDE_HEIGHT_PX = 540

#: Font size used when neither the run, the paragraph nor the placeholder kind
#: says one — PowerPoint would resolve this through the master's list styles,
#: which this renderer deliberately does not walk.
FALLBACK_SIZE_PT = 18.0
PLACEHOLDER_SIZE_PT = {
    PP_PLACEHOLDER.TITLE: 40.0,
    PP_PLACEHOLDER.CENTER_TITLE: 40.0,
    PP_PLACEHOLDER.VERTICAL_TITLE: 40.0,
    PP_PLACEHOLDER.SUBTITLE: 22.0,
    PP_PLACEHOLDER.SLIDE_NUMBER: 12.0,
    PP_PLACEHOLDER.FOOTER: 12.0,
    PP_PLACEHOLDER.HEADER: 12.0,
    PP_PLACEHOLDER.DATE: 12.0,
}
#: Placeholders whose paragraphs are bulleted unless the deck says otherwise.
BULLETED_PLACEHOLDERS = {
    PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.VERTICAL_BODY,
    PP_PLACEHOLDER.VERTICAL_OBJECT,
}

ALIGN_NAMES = {
    PP_ALIGN.LEFT: 'left', PP_ALIGN.CENTER: 'center', PP_ALIGN.RIGHT: 'right',
    PP_ALIGN.JUSTIFY: 'justify', PP_ALIGN.JUSTIFY_LOW: 'justify',
    PP_ALIGN.DISTRIBUTE: 'justify', PP_ALIGN.THAI_DISTRIBUTE: 'justify',
}
ANCHOR_NAMES = {
    MSO_ANCHOR.TOP: 'top', MSO_ANCHOR.MIDDLE: 'middle', MSO_ANCHOR.BOTTOM: 'bottom',
}

#: MSO_THEME_COLOR names -> the a:schemeClr token they stand for.
THEME_TOKENS = {
    'DARK_1': 'dk1', 'LIGHT_1': 'lt1', 'DARK_2': 'dk2', 'LIGHT_2': 'lt2',
    'TEXT_1': 'tx1', 'BACKGROUND_1': 'bg1', 'TEXT_2': 'tx2', 'BACKGROUND_2': 'bg2',
    'ACCENT_1': 'accent1', 'ACCENT_2': 'accent2', 'ACCENT_3': 'accent3',
    'ACCENT_4': 'accent4', 'ACCENT_5': 'accent5', 'ACCENT_6': 'accent6',
    'HYPERLINK': 'hlink', 'FOLLOWED_HYPERLINK': 'folHlink',
}
#: Tokens the master's p:clrMap re-points; the default mapping is the light one.
MAPPED_TOKENS = {'tx1': 'dk1', 'bg1': 'lt1', 'tx2': 'dk2', 'bg2': 'lt2'}

#: Text colour used where the deck leaves it to the master's list styles, picked
#: for contrast so a dark template does not render black-on-black.
LIGHT_TEXT = '#f5f5f5'
DARK_TEXT = '#1f2937'


class UnsupportedDeck(Exception):
    """The file cannot be rendered; the message is shown to the learner."""


# -- colour helpers ---------------------------------------------------------

def _clamp_byte(value):
    return max(0, min(255, int(round(value))))


def _apply_brightness(hex_color, brightness):
    """PowerPoint's lighten (>0, toward white) / darken (<0, toward black)."""
    if not brightness:
        return hex_color
    channels = [int(hex_color[index:index + 2], 16) for index in (1, 3, 5)]
    if brightness > 0:
        channels = [channel + (255 - channel) * brightness for channel in channels]
    else:
        channels = [channel * (1 + brightness) for channel in channels]
    return '#%02x%02x%02x' % tuple(_clamp_byte(channel) for channel in channels)


def _theme_colors(master):
    """The colour context for a master: theme slot hexes plus its colour map.

    A slide never names a theme slot directly — it names a token (tx1, bg1,
    accent1 …) that the master's p:clrMap points at a slot. A dark template
    inverts tx1/bg1, so ignoring the map renders its light body text as the
    near-black dk1 and its slides come out unreadable.
    """
    slots = {}
    try:
        root = etree.fromstring(master.part.part_related_by(RT.THEME).blob)
        scheme = root.find('.//' + qn('a:clrScheme'))
    except Exception:  # pragma: no cover - a deck without a readable theme
        scheme = None
    if scheme is not None:
        for slot in scheme:
            literal = _xml_color_hex(slot, {})
            if literal:
                slots[slot.tag.split('}')[-1]] = literal
    color_map = dict(MAPPED_TOKENS)
    mapping = master._element.find(qn('p:clrMap'))
    if mapping is not None:
        for token in MAPPED_TOKENS:
            if mapping.get(token):
                color_map[token] = str(mapping.get(token))
    return {'slots': slots, 'map': color_map}


def _scheme_hex(theme, token):
    """Hex for an a:schemeClr token, resolved through the master's colour map."""
    slots = (theme or {}).get('slots', {})
    color_map = (theme or {}).get('map', MAPPED_TOKENS)
    return slots.get(color_map.get(token, token))


def _xml_color_hex(element, theme):
    """First literal (or theme) colour found under an OOXML colour element."""
    if element is None:
        return None
    srgb = element.find('.//' + qn('a:srgbClr'))
    if srgb is not None and srgb.get('val'):
        return '#' + str(srgb.get('val')).lower()
    system = element.find('.//' + qn('a:sysClr'))
    if system is not None and system.get('lastClr'):
        return '#' + str(system.get('lastClr')).lower()
    scheme = element.find('.//' + qn('a:schemeClr'))
    if scheme is not None and scheme.get('val'):
        return _scheme_hex(theme, str(scheme.get('val')))
    return None


def _alpha_of(element):
    """Opacity 0-1 declared on an OOXML colour, or None when it is opaque."""
    if element is None:
        return None
    alpha = element.find('.//' + qn('a:alpha'))
    if alpha is None or not alpha.get('val'):
        return None
    try:
        value = int(alpha.get('val')) / 100000
    except (TypeError, ValueError):
        return None
    return value if 0 <= value < 1 else None


def _with_alpha(hex_color, alpha):
    """A CSS colour: plain hex when opaque, rgba() when not."""
    if alpha is None:
        return hex_color
    channels = tuple(int(hex_color[index:index + 2], 16) for index in (1, 3, 5))
    return 'rgba(%d, %d, %d, %s)' % (channels + (round(alpha, 3),))


def opaque_hex(css_color):
    """The '#rrggbb' behind a CSS colour, or None if it is translucent.

    Contrast decisions need something solid to measure; a translucent fill lets
    the slide behind it show through, so the caller falls back to that instead.
    """
    return css_color if css_color and len(css_color) == 7 and css_color[0] == '#' else None


def _color_hex(color, theme):
    """A CSS colour for a python-pptx ColorFormat, or None when it is unset."""
    if color is None:
        return None
    try:
        color_type = color.type
    except Exception:
        return None
    if color_type is None:
        return None
    resolved = None
    if color_type == MSO_COLOR_TYPE.RGB:
        try:
            resolved = '#' + str(color.rgb).lower()
        except Exception:
            resolved = None
    elif color_type == MSO_COLOR_TYPE.SCHEME:
        try:
            theme_color = color.theme_color
        except Exception:
            theme_color = None
        if isinstance(theme_color, MSO_THEME_COLOR):
            resolved = _scheme_hex(theme, THEME_TOKENS.get(theme_color.name, ''))
    if resolved is None:
        # Anything else (system, preset, HSL) still carries a literal value in
        # the XML often enough to be worth a look.
        resolved = _xml_color_hex(getattr(color, '_xFill', None), theme)
    if resolved is None:
        return None
    try:
        resolved = _apply_brightness(resolved, float(color.brightness or 0))
    except Exception:
        pass
    return _with_alpha(resolved, _alpha_of(getattr(color, '_xFill', None)))


def _fill_hex(fill, theme):
    """Solid colour for a shape/cell fill. Gradients collapse to their first stop."""
    if fill is None:
        return None
    try:
        fill_type = fill.type
    except Exception:
        return None
    if fill_type in (MSO_FILL.SOLID, MSO_FILL.PATTERNED):
        try:
            return _color_hex(fill.fore_color, theme)
        except Exception:
            return None
    if fill_type == MSO_FILL.GRADIENT:
        try:
            return _color_hex(fill.gradient_stops[0].color, theme)
        except Exception:
            return None
    return None


# -- geometry ---------------------------------------------------------------

class Frame:
    """Maps a shape's EMU coordinates onto the slide, scaled for group nesting."""

    def __init__(self, offset_x=0.0, offset_y=0.0, scale_x=1.0, scale_y=1.0):
        self.offset_x, self.offset_y = offset_x, offset_y
        self.scale_x, self.scale_y = scale_x, scale_y

    def px(self, left, top, width, height):
        return {
            'x': round((self.offset_x + (left or 0) * self.scale_x) / EMU_PER_PX, 2),
            'y': round((self.offset_y + (top or 0) * self.scale_y) / EMU_PER_PX, 2),
            'w': round((width or 0) * self.scale_x / EMU_PER_PX, 2),
            'h': round((height or 0) * self.scale_y / EMU_PER_PX, 2),
        }

    @property
    def text_scale(self):
        """Groups scale their text along with their geometry."""
        return (abs(self.scale_x) + abs(self.scale_y)) / 2 or 1.0

    def nested(self, group):
        """The frame a group's children are laid out in.

        A group defines its own coordinate space (chOff/chExt); children are
        placed in that space and stretched into the group's on-slide box.
        """
        fallback = Frame(
            self.offset_x + (group.left or 0) * self.scale_x,
            self.offset_y + (group.top or 0) * self.scale_y,
            self.scale_x, self.scale_y,
        )
        try:
            transform = group.element.find(qn('p:grpSpPr')).find(qn('a:xfrm'))
            child_x = int(transform.find(qn('a:chOff')).get('x'))
            child_y = int(transform.find(qn('a:chOff')).get('y'))
            child_cx = int(transform.find(qn('a:chExt')).get('cx')) or 1
            child_cy = int(transform.find(qn('a:chExt')).get('cy')) or 1
        except Exception:
            return fallback
        scale_x = self.scale_x * ((group.width or child_cx) / child_cx)
        scale_y = self.scale_y * ((group.height or child_cy) / child_cy)
        return Frame(
            self.offset_x + (group.left or 0) * self.scale_x - child_x * scale_x,
            self.offset_y + (group.top or 0) * self.scale_y - child_y * scale_y,
            scale_x, scale_y,
        )


# -- shape outlines ---------------------------------------------------------

#: prstGeom presets whose corners are rounded, and the guide that controls how
#: much. Everything else is drawn as its bounding box.
ROUNDED_PRESETS = {'roundRect', 'round1Rect', 'round2SameRect', 'round2DiagRect', 'snip1Rect'}
DEFAULT_ROUND_ADJ = 0.16667


def _svg_number(value):
    """A path coordinate SVG will accept.

    Path coordinates are raw EMU-scale integers, and "%g" would render those in
    scientific notation ("1.28e+07"), which SVG path data rejects outright.
    """
    rounded = round(float(value), 2)
    return '%d' % rounded if rounded == int(rounded) else '%.2f' % rounded


def _points_d(command, *points):
    return command + ' ' + ' '.join(
        _svg_number(coordinate) for point in points for coordinate in point
    )


def _point(element):
    """An a:pt as (x, y), or None when it is formula-driven rather than literal."""
    try:
        return float(element.get('x')), float(element.get('y'))
    except (TypeError, ValueError):
        return None


def _arc_segment(current, element):
    """An a:arcTo as an SVG elliptical-arc command, from the current point.

    OOXML gives the radii and a start/swing angle rather than an end point, so
    the end point has to be derived: the arc's centre sits one radius-vector
    back from where the pen already is.
    """
    try:
        radius_x, radius_y = float(element.get('wR')), float(element.get('hR'))
        start = math.radians(float(element.get('stAng')) / 60000)
        swing = math.radians(float(element.get('swAng')) / 60000)
    except (TypeError, ValueError):
        return None, current
    center = (current[0] - radius_x * math.cos(start), current[1] - radius_y * math.sin(start))
    end = (
        center[0] + radius_x * math.cos(start + swing),
        center[1] + radius_y * math.sin(start + swing),
    )
    command = 'A %s %s 0 %d %d %s %s' % (
        _svg_number(radius_x), _svg_number(radius_y),
        1 if abs(swing) > math.pi else 0, 1 if swing > 0 else 0,
        _svg_number(end[0]), _svg_number(end[1]),
    )
    return command, end


def _custom_path(shape):
    """A shape's a:custGeom as an SVG path, or None if it has none (or is odd).

    Freeform shapes carry the deck's design furniture — angled panels, cut
    corners, blobs — and drawing them as rectangles is what makes a rendered
    deck look nothing like the original.
    """
    try:
        geometry = shape._element.spPr.find(qn('a:custGeom'))
    except Exception:
        return None
    if geometry is None:
        return None
    path_list = geometry.find(qn('a:pathLst'))
    if path_list is None:
        return None
    commands = []
    view_w = view_h = 0
    for sub_path in path_list.findall(qn('a:path')):
        try:
            view_w = max(view_w, int(sub_path.get('w') or 0))
            view_h = max(view_h, int(sub_path.get('h') or 0))
        except (TypeError, ValueError):
            pass
        current = (0.0, 0.0)
        for step in sub_path:
            tag = step.tag.split('}')[-1]
            points = [_point(pt) for pt in step.findall(qn('a:pt'))]
            if any(point is None for point in points):
                return None  # formula-driven geometry: fall back to the box
            if tag == 'moveTo' and len(points) == 1:
                commands.append(_points_d('M', points[0]))
                current = points[0]
            elif tag == 'lnTo' and len(points) == 1:
                commands.append(_points_d('L', points[0]))
                current = points[0]
            elif tag == 'cubicBezTo' and len(points) == 3:
                commands.append(_points_d('C', *points))
                current = points[2]
            elif tag == 'quadBezTo' and len(points) == 2:
                commands.append(_points_d('Q', *points))
                current = points[1]
            elif tag == 'arcTo':
                command, current = _arc_segment(current, step)
                if command is None:
                    return None
                commands.append(command)
            elif tag == 'close':
                commands.append('Z')
    if not commands or not view_w or not view_h:
        return None
    return {'d': ' '.join(commands), 'vbW': view_w, 'vbH': view_h}


def _preset_radius(shape, box):
    """CSS corner radius for a rounded preset shape, or None."""
    try:
        preset = shape._element.spPr.find(qn('a:prstGeom'))
    except Exception:
        return None
    if preset is None:
        return None
    name = str(preset.get('prst') or '')
    if name == 'ellipse':
        return '50%'
    if name not in ROUNDED_PRESETS:
        return None
    adjust = DEFAULT_ROUND_ADJ
    guide = preset.find('.//' + qn('a:gd'))
    if guide is not None and str(guide.get('fmla') or '').startswith('val '):
        try:
            adjust = int(str(guide.get('fmla'))[4:]) / 100000
        except (TypeError, ValueError):
            adjust = DEFAULT_ROUND_ADJ
    return '%gpx' % round(min(box['w'], box['h']) * min(max(adjust, 0), 0.5), 2)


# -- text -------------------------------------------------------------------

def _bullet_for(paragraph, default_bulleted):
    """The bullet glyph for a paragraph, or None when it has no bullet.

    Only the paragraph's own properties are read: resolving inherited bullets
    would mean walking the master's list styles, so a deck that leaves bullets
    entirely to its master renders as clean text rather than wrong text.
    """
    properties = paragraph._pPr
    if properties is not None:
        if properties.find(qn('a:buNone')) is not None:
            return None
        char = properties.find(qn('a:buChar'))
        if char is not None:
            return str(char.get('char') or '•')
        if properties.find(qn('a:buAutoNum')) is not None:
            return '•'
    return '•' if default_bulleted else None


def _font_scale(text_frame):
    """PowerPoint's "shrink text on overflow" factor, as stored in the file."""
    try:
        body = text_frame._txBody.find(qn('a:bodyPr'))
        autofit = body.find(qn('a:normAutofit')) if body is not None else None
        if autofit is not None and autofit.get('fontScale'):
            return int(autofit.get('fontScale')) / 100000
    except Exception:
        pass
    return 1.0


def _run_size_pt(run, paragraph, fallback_pt):
    for font in (run.font, paragraph.font):
        try:
            if font.size is not None:
                return font.size.pt
        except Exception:
            continue
    return fallback_pt


def _paragraph_pieces(paragraph):
    """Runs in document order, with soft line breaks and field text kept.

    `paragraph.runs` sees only <a:r>, so a line break inside a title (<a:br/>)
    and an auto-numbered slide number (<a:fld>) would both vanish. Walking the
    children keeps them.
    """
    for child in paragraph._p:
        tag = child.tag.split('}')[-1]
        if tag == 'r':
            yield _Run(child, paragraph)
        elif tag == 'br':
            yield '\n'
        elif tag == 'fld':
            text_element = child.find(qn('a:t'))
            if text_element is not None and text_element.text:
                yield text_element.text


def _normalise_run_text(text):
    """Soft breaks and stray carriage returns become plain newlines.

    The viewer renders run text with `white-space: pre-wrap`, so a newline
    here is the line break the author typed.
    """
    return str(text).replace('\x0b', '\n').replace('\r\n', '\n').replace('\r', '\n')


def _paragraph_model(paragraph, theme, fallback_pt, default_bulleted, scale):
    runs = []
    for piece in _paragraph_pieces(paragraph):
        if isinstance(piece, str):
            # A break or field belongs to the run it follows, so it inherits its
            # styling instead of introducing a differently-sized fragment.
            if runs:
                runs[-1]['text'] += _normalise_run_text(piece)
            continue
        if not piece.text:
            continue
        runs.append({
            'text': _normalise_run_text(piece.text),
            'sizePx': round(_run_size_pt(piece, paragraph, fallback_pt) * PT_TO_PX * scale, 2),
            'bold': bool(piece.font.bold),
            'italic': bool(piece.font.italic),
            'underline': bool(piece.font.underline),
            'color': _color_hex(piece.font.color, theme),
            'font': piece.font.name or None,
        })
    if not runs:
        return None
    line_spacing, line_height_px = None, None
    try:
        spacing = paragraph.line_spacing
        if isinstance(spacing, Length):
            line_height_px = round(spacing.pt * PT_TO_PX * scale, 2)
        elif isinstance(spacing, (int, float)):
            line_spacing = round(float(spacing), 3)
    except Exception:
        line_spacing, line_height_px = None, None
    return {
        'align': ALIGN_NAMES.get(paragraph.alignment),
        'level': int(paragraph.level or 0),
        'bullet': _bullet_for(paragraph, default_bulleted),
        'lineSpacing': line_spacing,
        'lineHeightPx': line_height_px,
        'runs': runs,
    }


def _text_model(text_frame, theme, fallback_pt, default_bulleted, scale):
    font_scale = _font_scale(text_frame) * scale
    paragraphs = [
        model for model in (
            _paragraph_model(paragraph, theme, fallback_pt, default_bulleted, font_scale)
            for paragraph in text_frame.paragraphs
        ) if model
    ]
    try:
        anchor = ANCHOR_NAMES.get(text_frame.vertical_anchor, 'top')
    except Exception:
        anchor = 'top'
    return {'valign': anchor, 'paragraphs': paragraphs}


# -- shapes -----------------------------------------------------------------

def _relative_luminance(hex_color):
    """0 (black) to 1 (white), the sRGB luminance used for contrast decisions."""
    channels = []
    for index in (1, 3, 5):
        channel = int(hex_color[index:index + 2], 16) / 255
        channels.append(channel / 12.92 if channel <= 0.04045 else ((channel + 0.055) / 1.055) ** 2.4)
    return 0.2126 * channels[0] + 0.7152 * channels[1] + 0.0722 * channels[2]


def _contrast_text(background_hex):
    """Readable text colour for whatever the text is drawn on."""
    if not background_hex or len(background_hex) != 7:
        return DARK_TEXT
    try:
        return DARK_TEXT if _relative_luminance(background_hex) > 0.4 else LIGHT_TEXT
    except ValueError:
        return DARK_TEXT


def _is_readable_on(text_hex, surface_hex):
    """Rough "can this be read" test, used only to reject invisible text."""
    if not text_hex or not surface_hex or len(text_hex) != 7 or len(surface_hex) != 7:
        return True
    try:
        return abs(_relative_luminance(text_hex) - _relative_luminance(surface_hex)) > 0.12
    except ValueError:
        return True


def _style_font_color(shape, theme):
    """The colour a styled shape gives its text (p:style/a:fontRef).

    This is PowerPoint's own answer for a shape drawn from a theme style — a
    dark chip whose caption is white, say — and it beats any guess of ours.
    """
    try:
        style = shape._element.find(qn('p:style'))
        font_ref = style.find(qn('a:fontRef')) if style is not None else None
    except Exception:
        return None
    return _xml_color_hex(font_ref, theme) if font_ref is not None else None


def _default_text_color(shape, fill, theme, background_hex):
    """Colour for runs that name none of their own.

    PowerPoint resolves this through the master's list styles, which this
    renderer does not walk. The next-best sources are the shape's own style and
    the theme's body-text slot; either is used only if the result can actually
    be read against what it sits on, so no deck renders as invisible text.
    """
    surface = opaque_hex(fill) or background_hex
    for candidate in (_style_font_color(shape, theme), _scheme_hex(theme, 'tx1')):
        if candidate and _is_readable_on(candidate, surface):
            return candidate
    return _contrast_text(surface)


def _placeholder_type(shape):
    try:
        return shape.placeholder_format.type if shape.is_placeholder else None
    except Exception:
        return None


def _rotation(shape):
    try:
        rotation = float(shape.rotation or 0)
    except Exception:
        return None
    return round(rotation, 2) if rotation else None


def _image_source(image, context):
    """Write an embedded image out once and return the URL it is served from."""
    try:
        blob = image.blob
    except Exception:
        return None
    digest = hashlib.sha1(blob).hexdigest()[:16]
    assets = context['assets']
    if digest in assets:
        return assets[digest]
    extension = (getattr(image, 'ext', None) or 'png').lower().lstrip('.')
    name = 'img-%s.%s' % (digest, extension)
    target = context['asset_dir'] / name
    target.parent.mkdir(parents=True, exist_ok=True)
    if not target.exists():
        target.write_bytes(blob)
    assets[digest] = '%s/%s' % (context['url_prefix'], name)
    return assets[digest]


def _table_model(shape, box, theme, scale):
    table = shape.table
    total_width = sum(column.width or 0 for column in table.columns) or 1
    total_height = sum(row.height or 0 for row in table.rows) or 1
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            if getattr(cell, 'is_spanned', False):
                continue
            cells.append({
                'colSpan': int(getattr(cell, 'span_width', 1) or 1),
                'rowSpan': int(getattr(cell, 'span_height', 1) or 1),
                'fill': _fill_hex(cell.fill, theme),
                **_text_model(cell.text_frame, theme, 12.0, False, scale),
            })
        rows.append(cells)
    return {
        'kind': 'table',
        **box,
        'colFractions': [round((column.width or 0) / total_width, 4) for column in table.columns],
        'rowFractions': [round((row.height or 0) / total_height, 4) for row in table.rows],
        'rows': rows,
    }


UNSUPPORTED_LABELS = {
    MSO_SHAPE_TYPE.CHART: 'Chart',
    MSO_SHAPE_TYPE.DIAGRAM: 'Diagram',
    MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: 'Embedded object',
    MSO_SHAPE_TYPE.LINKED_OLE_OBJECT: 'Linked object',
    MSO_SHAPE_TYPE.MEDIA: 'Media clip',
}


def _shape_models(shape, frame, context):
    """One shape — recursively, for groups — as zero or more model entries."""
    theme = context['theme']
    box = frame.px(shape.left, shape.top, shape.width, shape.height)
    shape_type = getattr(shape, 'shape_type', None)

    if shape_type == MSO_SHAPE_TYPE.GROUP:
        nested = frame.nested(shape)
        models = []
        for child in shape.shapes:
            models.extend(_shape_models(child, nested, context))
        return models

    if shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE):
        source = _image_source(shape.image, context)
        if not source:
            return []
        return [{
            'kind': 'image', **box, 'src': source,
            'alt': shape.name or 'Slide image', 'rotation': _rotation(shape),
        }]

    if getattr(shape, 'has_table', False):
        return [_table_model(shape, box, theme, frame.text_scale)]

    if shape_type in UNSUPPORTED_LABELS:
        return [{'kind': 'note', **box, 'label': UNSUPPORTED_LABELS[shape_type]}]

    fill = _fill_hex(getattr(shape, 'fill', None), theme)
    line, line_width_pt = None, None
    try:
        line = _color_hex(shape.line.color, theme)
        line_width_pt = shape.line.width.pt if shape.line.width else None
    except Exception:
        line, line_width_pt = None, None
    has_text = bool(getattr(shape, 'has_text_frame', False))
    if not has_text and not fill and not line:
        return []

    outline = _custom_path(shape)
    model = {
        'kind': 'text' if has_text else 'shape',
        **box,
        'path': outline,
        'radius': None if outline else _preset_radius(shape, box),
        'fill': fill,
        'line': line,
        'lineWidthPx': round(line_width_pt * PT_TO_PX, 2) if line_width_pt else None,
        'rotation': _rotation(shape),
        'valign': 'top',
        # Runs that name no colour of their own inherit this.
        'defaultTextColor': _default_text_color(shape, fill, theme, context['background_hex']),
        'paragraphs': [],
    }
    if has_text:
        placeholder = _placeholder_type(shape)
        model.update(_text_model(
            shape.text_frame, theme,
            PLACEHOLDER_SIZE_PT.get(placeholder, FALLBACK_SIZE_PT),
            placeholder in BULLETED_PLACEHOLDERS,
            frame.text_scale,
        ))
        if not model['paragraphs'] and not fill and not line:
            return []
    return [model]


# -- backgrounds and inherited (layout / master) graphics -------------------

def _background(element, part, theme, context):
    """{'color': …} or {'image': …} for a slide/layout/master's own background."""
    common = element.find(qn('p:cSld'))
    background = common.find(qn('p:bg')) if common is not None else None
    if background is None:
        return None
    blip = background.find('.//' + qn('a:blip'))
    if blip is not None and blip.get(qn('r:embed')):
        try:
            image_part = part.related_part(blip.get(qn('r:embed')))
            source = _image_source(image_part, context)
            if source:
                return {'image': source}
        except Exception:
            pass
    color = _xml_color_hex(background, theme)
    return {'color': color} if color else None


def _follows_master_graphics(element):
    return element.get('showMasterSp') not in ('0', 'false')


def _inherited_shapes(slide, context):
    """The template's own graphics: master then layout, placeholders excluded.

    A designed deck keeps its furniture (background bars, logos, rules) on the
    master and its layouts, so skipping these renders every slide as bare text
    on white. Placeholders are skipped — on a master or layout they are empty
    prompts that the slide itself fills in.
    """
    layout = slide.slide_layout
    master = layout.slide_master
    models = []
    if not _follows_master_graphics(slide._element):
        return models
    if _follows_master_graphics(layout._element):
        for shape in master.shapes:
            if not shape.is_placeholder:
                models.extend(_shape_models(shape, Frame(), context))
    for shape in layout.shapes:
        if not shape.is_placeholder:
            models.extend(_shape_models(shape, Frame(), context))
    return models


# -- the deck ---------------------------------------------------------------

def _deck_model(source, asset_dir, url_prefix):
    presentation = Presentation(str(source))
    assets = {}
    slides = []
    for number, slide in enumerate(presentation.slides, start=1):
        layout = slide.slide_layout
        master = layout.slide_master
        context = {
            'theme': _theme_colors(master),
            'assets': assets,
            'asset_dir': asset_dir,
            'url_prefix': url_prefix,
            'background_hex': '#ffffff',
        }
        theme = context['theme']
        background = (
            _background(slide._element, slide.part, theme, context)
            or _background(layout._element, layout.part, theme, context)
            or _background(master._element, master.part, theme, context)
        )
        context['background_hex'] = (background or {}).get('color') or '#ffffff'
        shapes = _inherited_shapes(slide, context)
        for shape in slide.shapes:
            shapes.extend(_shape_models(shape, Frame(), context))
        try:
            notes = (slide.notes_slide.notes_text_frame.text or '').strip() if slide.has_notes_slide else ''
        except Exception:
            notes = ''
        slides.append({
            'number': number,
            'layout': layout.name or '',
            'background': background,
            'shapes': shapes,
            'notes': notes,
        })
    return {
        'slideWidthPx': round((presentation.slide_width or 0) / EMU_PER_PX) or DEFAULT_SLIDE_WIDTH_PX,
        'slideHeightPx': round((presentation.slide_height or 0) / EMU_PER_PX) or DEFAULT_SLIDE_HEIGHT_PX,
        'slideCount': len(slides),
        'slides': slides,
    }


def render_root():
    from curriculum_api.views import COMPONENT_UPLOAD_ROOT
    return Path(settings.MEDIA_ROOT) / COMPONENT_UPLOAD_ROOT / RENDER_DIR_NAME


def _stamp(source):
    stat = source.stat()
    return '%s:%s:%s' % (MODEL_VERSION, stat.st_size, stat.st_mtime_ns)


def render_cache_key(relative_path):
    """The cache directory name for an upload path."""
    return hashlib.sha1(relative_path.encode('utf-8')).hexdigest()[:20]


def cached_deck_model(relative_path):
    """The last render stored for this upload path, or None.

    Read-only and stamp-blind, unlike ``render_uploaded_deck``, which exists to
    keep a render in step with its source: that one compares ``_stamp(source)``
    and deletes the cache when the file has changed underneath it. Neither is
    possible once the source is gone — there is no stat to take and nothing to
    compare — and deleting the only surviving copy of a deck would be the worst
    of the available options.

    So this is deliberately the lenient reader, used only when the upload has
    vanished. Uploads live under MEDIA_ROOT, which is not in version control and
    not backed up, so a file disappearing while its render survives is a real
    state and not a theoretical one.
    """
    manifest = render_root() / render_cache_key(relative_path) / 'deck.json'
    if not manifest.is_file():
        return None
    try:
        cached = json.loads(manifest.read_text(encoding='utf-8'))
    except (OSError, ValueError):
        logger.warning('Unreadable slide render cache at %s', manifest)
        return None
    deck = cached.get('deck')
    return deck if isinstance(deck, dict) else None


def render_uploaded_deck(relative_path, source):
    """The slide model for an uploaded deck, rendering it if not already cached.

    `relative_path` is the storage-relative path of the upload and keys the
    cache. Raises UnsupportedDeck, whose message is learner-facing, when the
    file is not an OOXML deck or cannot be parsed.
    """
    source = Path(source)
    suffix = source.suffix.lower()
    if suffix in LEGACY_SUFFIXES:
        raise UnsupportedDeck(
            'This deck uses the older PowerPoint 97-2003 format, which cannot be shown '
            'inline. Re-save it as .pptx — or as a PDF — and upload it again.'
        )
    if suffix not in OOXML_SUFFIXES:
        raise UnsupportedDeck('This file is not a PowerPoint deck that can be shown inline.')

    key = hashlib.sha1(relative_path.encode('utf-8')).hexdigest()[:20]
    cache_dir = render_root() / key
    manifest = cache_dir / 'deck.json'
    url_prefix = '/curriculum_api/curriculum/uploads/%s/%s' % (RENDER_DIR_NAME, key)
    stamp = _stamp(source)

    if manifest.exists():
        try:
            cached = json.loads(manifest.read_text(encoding='utf-8'))
            if cached.get('stamp') == stamp:
                return cached['deck']
        except Exception:
            logger.warning('Discarding unreadable slide render cache at %s', manifest)
    # Any surviving render belongs to a replaced file; its images are dead weight.
    if cache_dir.exists():
        shutil.rmtree(cache_dir, ignore_errors=True)
    cache_dir.mkdir(parents=True, exist_ok=True)

    try:
        deck = _deck_model(source, cache_dir, url_prefix)
    except Exception as error:
        logger.exception('Could not render slide deck %s', relative_path)
        raise UnsupportedDeck(
            'This deck could not be read for inline display. Open or download it instead.'
        ) from error

    manifest.write_text(
        json.dumps({'stamp': stamp, 'source': relative_path, 'deck': deck}),
        encoding='utf-8',
    )
    return deck


def is_renderable_deck_path(path):
    """True when a resource path looks like a deck this module can render."""
    cleaned = str(path or '').split('?')[0].split('#')[0]
    return Path(cleaned).suffix.lower() in OOXML_SUFFIXES
