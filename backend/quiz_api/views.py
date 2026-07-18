import json
import csv
import html
import io
import base64
import logging
import mimetypes
import re
import zipfile
from pathlib import Path
from posixpath import normpath
from xml.etree import ElementTree

from django.conf import settings
from django.db import DatabaseError, connection, transaction
from django.db.models import Max
from django.http import FileResponse, Http404, HttpResponse, JsonResponse
from django.utils import timezone
from django.views.decorators.clickjacking import xframe_options_sameorigin
from django.views.decorators.csrf import csrf_exempt
from django.views.decorators.http import require_http_methods

from .models import QuizAnswer, QuizPackage, QuizQuestion
from learner_api.models import ActiveUser

logger = logging.getLogger(__name__)

AI_GENERATION_MAX_FILE_SIZE = 50 * 1024 * 1024
AI_GENERATION_MAX_CONTENT_CHARS = 60000
AI_GENERATION_MAX_OCR_PDF_PAGES = 24
AI_GENERATION_OCR_PAGE_BATCH_SIZE = 4
AI_GENERATION_TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".xml", ".html", ".htm", ".json"}
AI_GENERATION_QUESTION_TYPES = {
    "single_choice",
    "multiple_choice",
    "true_false",
    "matching",
    "image_matching",
    "keywords",
    "fill_gap",
    "ordering",
}
QUIZ_STATUSES = {"draft", "published", "pending", "validating", "trash", "private"}
ASSESSMENT_TYPES = {"quiz", "checkpoint"}


def _ensure_quiz_assessment_type_column():
    with connection.cursor() as cursor:
        cursor.execute(
            """
            do $$
            begin
              if exists (
                select 1
                from information_schema.columns
                where table_schema = 'curriculum'
                  and table_name = 'quizzes'
                  and column_name = 'programme_id'
                  and data_type <> 'character varying'
              ) then
                alter table curriculum.quizzes
                alter column programme_id type varchar(128)
                using programme_id::varchar;
              end if;
            end $$;
            """
        )
        cursor.execute(
            """
            with matches as (
              select
                q.id as quiz_id,
                mam.programme_id,
                row_number() over (
                  partition by q.id
                  order by
                    case
                      when mam.imported_from_training_plan_id = q.programme_id then 0
                      when lower(mam.title) = lower(q.module) then 1
                      else 2
                    end,
                    mam.updated_at desc nulls last,
                    mam.created_at desc nulls last
                ) as match_rank
              from curriculum.quizzes q
              left join curriculum."Training_plan" tp
                on q.programme_id ~ '^[0-9]+$'
               and tp.id = q.programme_id::integer
              join curriculum.module_authoring_modules mam
                on coalesce(trim(mam.programme_id), '') <> ''
               and (
                 mam.imported_from_training_plan_id = q.programme_id
                 or (
                   (lower(mam.title) = lower(q.module) or lower(q.module) like ('%%' || lower(mam.title) || '%%'))
                   and (
                     lower(mam.programme_name) = lower(q.programme)
                     or lower(mam.programme_id) = lower(q.programme)
                     or lower(mam.programme_name) = lower(tp."Program")
                     or lower(mam.programme_id) = lower(tp."Program")
                   )
                 )
                 or lower(mam.programme_name) = lower(tp."Program")
                 or lower(mam.programme_id) = lower(tp."Program")
               )
              where q.programme_id ~ '^[0-9]+$'
            )
            update curriculum.quizzes q
            set programme_id = matches.programme_id
            from matches
            where q.id = matches.quiz_id
              and matches.match_rank = 1
            """
        )
        cursor.execute(
            """
            update curriculum.quizzes
            set programme_id = programme
            where programme_id ~ '^[0-9]+$'
              and coalesce(trim(programme), '') <> ''
              and programme !~ '^[0-9]+$'
            """
        )
        cursor.execute(
            """
            alter table curriculum.quizzes
            add column if not exists assessment_type varchar(40) not null default 'quiz'
            """
        )
        cursor.execute(
            """
            alter table curriculum.quizzes
            add column if not exists week_id varchar(128) not null default ''
            """
        )


def _clean_assessment_type(value, default="quiz"):
    assessment_type = str(value or default).strip().lower()
    return assessment_type if assessment_type in ASSESSMENT_TYPES else default


def _week_number_from_value(value):
    text = str(value or "").strip()
    if not text:
        return None
    if text.isdigit():
        return int(text)
    week_match = re.search(r"\bweek\s*(\d+)\b", text, flags=re.IGNORECASE)
    if week_match:
        return int(week_match.group(1))
    compact_match = re.search(r"(?:^|[-_\s])w(?:eek)?[-_\s]*(\d+)(?:$|[-_\s])", text, flags=re.IGNORECASE)
    if compact_match:
        return int(compact_match.group(1))
    return None


def _build_week_id(programme_id, week_value):
    week_number = _week_number_from_value(week_value)
    if not programme_id or not week_number:
        return ""
    return f"week-training-module-{programme_id}-{week_number}"


def _is_int_like(value):
    return str(value or "").strip().isdigit()


def _format_size(size):
    if size >= 1024 * 1024:
        return f"{size / (1024 * 1024):.1f} MB"
    if size >= 1024:
        return f"{round(size / 1024)} KB"
    return f"{size} B"


def _serialize_quiz(quiz):
    return {
        "id": quiz.id,
        "title": quiz.title,
        "programmeId": quiz.programme_id,
        "module": quiz.module,
        "programme": quiz.programme,
        "version": quiz.version,
        "questions": quiz.questions,
        "defaultQuestionType": quiz.default_question_type,
        "assessmentType": quiz.assessment_type,
        "weekId": quiz.week_id,
        "status": quiz.status,
        "packageType": quiz.package_type,
        "fileName": quiz.file_name,
        "fileSize": quiz.file_size,
        "displaySize": _format_size(quiz.file_size),
        "schemaValid": quiz.schema_valid,
        "validationMessage": quiz.validation_message,
        "mappedComponents": quiz.mapped_components,
        "author": quiz.author,
        "linkedCourses": quiz.linked_courses,
        "shortDescription": quiz.short_description,
        "lessonContent": quiz.lesson_content,
        "duration": quiz.duration,
        "timeUnit": quiz.time_unit,
        "quizStyle": quiz.quiz_style,
        "randomizeQuestions": quiz.randomize_questions,
        "randomizeAnswers": quiz.randomize_answers,
        "showCorrectAnswer": quiz.show_correct_answer,
        "attemptHistory": quiz.attempt_history,
        "retakeAfterPass": quiz.retake_after_pass,
        "limitAttempts": quiz.limit_attempts,
        "passingGrade": quiz.passing_grade,
        "retakePointsCut": quiz.retake_points_cut,
        "publishedAt": quiz.published_at.isoformat() if quiz.published_at else None,
        "createdAt": quiz.created_at.isoformat(),
        "updatedAt": quiz.updated_at.isoformat(),
    }


def _clean_quiz_status(value, default="draft"):
    status = str(value or default).strip().lower()
    return status if status in QUIZ_STATUSES else default


def _extract_text_for_ai(uploaded_file):
    if not uploaded_file:
        return ""
    if uploaded_file.size > AI_GENERATION_MAX_FILE_SIZE:
        raise ValueError("File is too large for question generation. Upload a file under 50 MB.")

    extension = Path(uploaded_file.name).suffix.lower()
    uploaded_file.seek(0)

    if extension in AI_GENERATION_TEXT_EXTENSIONS:
        return uploaded_file.read().decode("utf-8-sig", errors="replace")

    if extension in {".xlsx", ".xlsm"}:
        from openpyxl import load_workbook

        workbook = load_workbook(uploaded_file, read_only=True, data_only=True)
        lines = []
        for sheet in workbook.worksheets[:3]:
            lines.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if values:
                    lines.append(" | ".join(values))
        return "\n".join(lines)

    if extension in {".pptx", ".pptm"}:
        from pptx import Presentation

        presentation = Presentation(uploaded_file)
        lines = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = re.sub(r"\s+", " ", shape.text).strip()
                    if text:
                        slide_lines.append(text)
            if slide_lines:
                lines.append(f"Slide {slide_index}: " + " | ".join(slide_lines))
        return "\n".join(lines)

    if extension == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(uploaded_file)
        lines = []
        for page_index, page in enumerate(reader.pages[:30], start=1):
            text = page.extract_text() or ""
            text = re.sub(r"\s+", " ", text).strip()
            if text:
                lines.append(f"Page {page_index}: {text}")
        text = "\n".join(lines)
        if text:
            return text
        return _extract_pdf_text_with_vision(uploaded_file)

    if extension == ".docx":
        from docx import Document

        document = Document(uploaded_file)
        lines = []
        for paragraph in document.paragraphs:
            text = re.sub(r"\s+", " ", paragraph.text).strip()
            if text:
                lines.append(text)
        for table in document.tables:
            for row in table.rows:
                values = [re.sub(r"\s+", " ", cell.text).strip() for cell in row.cells if cell.text.strip()]
                if values:
                    lines.append(" | ".join(values))
        return "\n".join(lines)

    if extension == ".zip":
        lines = []
        with zipfile.ZipFile(uploaded_file) as archive:
            text_members = [
                member for member in archive.infolist()
                if not member.is_dir()
                and Path(member.filename).suffix.lower() in AI_GENERATION_TEXT_EXTENSIONS
                and member.file_size <= 512 * 1024
            ]
            for member in text_members[:40]:
                with archive.open(member) as file:
                    text = file.read().decode("utf-8-sig", errors="replace")
                    text = re.sub(r"\s+", " ", text).strip()
                    if text:
                        lines.append(f"File {member.filename}: {text}")
        return "\n".join(lines)

    raise ValueError("Unsupported AI source file. Upload TXT, CSV, XML, HTML, XLSX, XLSM, PPTX, PPTM, PDF, DOCX or SCORM ZIP.")


def _extract_pdf_text_with_vision(uploaded_file):
    try:
        import fitz
        from openai import OpenAI
    except ImportError:
        return ""

    uploaded_file.seek(0)
    try:
        document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
    except Exception:
        uploaded_file.seek(0)
        return ""

    if not document.page_count:
        return ""

    client = OpenAI(api_key=settings.OPENAI_API_KEY)
    extracted_pages = []
    total_pages = min(document.page_count, AI_GENERATION_MAX_OCR_PDF_PAGES)

    for batch_start in range(0, total_pages, AI_GENERATION_OCR_PAGE_BATCH_SIZE):
        content = [{
            "type": "input_text",
            "text": (
                "Extract all readable text from these PDF page images for LMS quiz generation. "
                "Preserve slide/page headings and bullet meaning. Return plain text only, grouped by page. "
                "Do not summarise and do not add external knowledge."
            ),
        }]
        for page_index in range(batch_start, min(batch_start + AI_GENERATION_OCR_PAGE_BATCH_SIZE, total_pages)):
            page = document.load_page(page_index)
            pixmap = page.get_pixmap(dpi=140, alpha=False)
            image_data = base64.b64encode(pixmap.tobytes("png")).decode("ascii")
            content.append({
                "type": "input_image",
                "image_url": f"data:image/png;base64,{image_data}",
            })

        try:
            response = client.responses.create(
                model=settings.OPENAI_MODEL,
                input=[{"role": "user", "content": content}],
            )
            batch_text = response.output_text.strip()
            if batch_text:
                extracted_pages.append(batch_text)
        except Exception:
            continue

    uploaded_file.seek(0)
    return "\n\n".join(extracted_pages)


def _extract_text_from_ai_files_with_report(files):
    extracted_chunks = []
    unreadable_files = []
    for uploaded_file in sorted(files[:10], key=lambda file: file.name.lower()):
        try:
            extracted = _extract_text_for_ai(uploaded_file).strip()
        except ValueError:
            # File-level validation problems (too large, unsupported type) must
            # surface to the caller as a 400, not be swallowed here.
            raise
        except Exception:
            # A single unreadable/corrupt file, or a missing optional extraction
            # dependency, should mark that file unreadable rather than 500 the
            # whole generation request.
            logger.exception("AI source extraction failed for %s", getattr(uploaded_file, "name", "?"))
            extracted = ""
        if extracted:
            extracted_chunks.append((uploaded_file.name, extracted))
        else:
            unreadable_files.append(uploaded_file.name)

    if not extracted_chunks:
        return "", [], unreadable_files

    per_file_budget = max(6000, AI_GENERATION_MAX_CONTENT_CHARS // len(extracted_chunks))
    chunks = [
        f"Source file: {name}\n{text[:per_file_budget]}"
        for name, text in extracted_chunks
    ]
    return "\n\n".join(chunks)[:AI_GENERATION_MAX_CONTENT_CHARS], [name for name, _ in extracted_chunks], unreadable_files


def _extract_text_from_ai_files(files):
    source_text, _, _ = _extract_text_from_ai_files_with_report(files)
    return source_text


def _source_file_names_from_text(source_text):
    return [name.strip() for name in re.findall(r"^Source file:\s*(.+)$", source_text, flags=re.MULTILINE)]


def _question_allocation_instruction(source_files, question_count):
    if not source_files:
        return "No uploaded source-file split was detected; distribute questions across the distinct themes in the source content."

    readable_files = list(dict.fromkeys(source_files))
    if question_count < len(readable_files):
        selected = readable_files[:question_count]
        return (
            "Required source coverage: create one question from each of these files and do not repeat a file: "
            + "; ".join(selected)
            + "."
        )

    base_count = question_count // len(readable_files)
    remainder = question_count % len(readable_files)
    allocations = []
    for index, name in enumerate(readable_files):
        extra = 1 if index >= len(readable_files) - remainder and remainder else 0
        allocations.append((name, base_count + extra))

    return (
        "Required source-file allocation: "
        + "; ".join(f"{name} = exactly {count} questions" for name, count in allocations)
        + ". This allocation is mandatory; do not overrepresent one file."
    )


def _normalise_generated_questions(data, default_type="single_choice"):
    raw_questions = data.get("questions", []) if isinstance(data, dict) else []
    normalised = []
    for question_index, item in enumerate(raw_questions):
        if not isinstance(item, dict):
            continue
        text = str(item.get("question") or item.get("text") or "").strip()
        if not text:
            continue
        question_type = _normalise_question_type(item.get("type") or default_type)
        raw_options = item.get("options") if isinstance(item.get("options"), list) else []
        options = [str(option).strip() for option in raw_options if str(option).strip()]
        correct_answer = str(item.get("correct_answer") or item.get("correctAnswer") or "").strip()
        explanation = str(item.get("explanation") or item.get("feedback") or "").strip()

        if question_type == "true_false" and len(options) < 2:
            options = ["True", "False"]
        if question_type in {"matching", "image_matching"} and not options:
            pair_source = correct_answer or explanation
            if "|" in pair_source:
                pair_source = pair_source.split("|", 1)[1]
            options = [part.strip() for part in re.split(r"\s*(?:;|\n)\s*", pair_source) if part.strip()]
        if question_type == "keywords" and not options:
            keyword_source = correct_answer or explanation
            if "|" in keyword_source:
                keyword_source = keyword_source.split("|", 1)[1]
            options = [part.strip() for part in re.split(r"\s*(?:,|;|\n|\|)\s*", keyword_source) if part.strip()]
        if question_type == "fill_gap" and not options and correct_answer:
            options = [correct_answer]
        if question_type == "ordering" and not correct_answer:
            correct_answer = ", ".join(chr(65 + index) for index in range(len(options)))
        if question_type in {"single_choice", "multiple_choice", "true_false"} and not options:
            continue

        correct_indexes = set(range(len(options))) if question_type in {"matching", "image_matching", "keywords", "ordering"} else {0}
        if correct_answer:
            correct_indexes = set()
            correct_tokens = [token.strip() for token in re.split(r"[,;/]", correct_answer) if token.strip()]
            for token in correct_tokens or [correct_answer]:
                lower_correct = token.lower()
                for index, option in enumerate(options):
                    if option.lower() == lower_correct:
                        correct_indexes.add(index)
                letter_index = _correct_answer_index(token)
                if 0 <= letter_index < len(options):
                    correct_indexes.add(letter_index)
            if question_type in {"matching", "image_matching", "keywords", "ordering"} and not correct_indexes:
                correct_indexes = set(range(len(options)))
            if not correct_indexes:
                correct_indexes = {0}

        answers = [
            {
                "id": -((question_index + 1) * 100 + answer_index + 1),
                "text": option,
                "isCorrect": answer_index in correct_indexes,
            }
            for answer_index, option in enumerate(options)
        ]
        normalised.append({
            "id": -(question_index + 1),
            "text": text,
            "questionType": question_type,
            "explanation": explanation,
            "answers": answers,
        })
    return normalised


def _program_candidates(programme):
    if not programme:
        return []

    cleaned = re.sub(r"\s+", " ", programme).strip()
    candidates = [cleaned]
    level_match = re.search(r"(?:level|lvl|l)\s*(\d+)", cleaned, flags=re.IGNORECASE)
    words = [word for word in re.findall(r"[A-Za-z]+", cleaned) if word.lower() not in {"level", "lvl"}]
    if words and level_match:
        initials = "".join(word[0].upper() for word in words[:2])
        candidates.append(f"{initials} L{level_match.group(1)}")
    return list(dict.fromkeys(candidates))


def _match_training_plan_id(programme, module, title):
    candidates = _program_candidates(programme)
    module_text = module or ""
    title_text = title or ""

    try:
        with connection.cursor() as cursor:
            if candidates:
                program_clauses = []
                params = []
                for candidate in candidates:
                    program_clauses.append('(lower("Program") = lower(%s) or lower("Program") like %s)')
                    params.extend([candidate, f"%{candidate.lower()}%"])

                cursor.execute(
                    f"""
                    select id
                    from curriculum."Training_plan"
                    where {" or ".join(program_clauses)}
                    order by
                      case
                        when lower(module_name) = lower(%s) then 0
                        when lower(%s) like ('%%' || lower(module_name) || '%%') then 1
                        when lower(%s) like ('%%' || lower(module_name) || '%%') then 2
                        else 3
                      end,
                      length(module_name) desc,
                      id desc
                    limit 1
                    """,
                    [*params, module_text, module_text, title_text],
                )
                row = cursor.fetchone()
                if row:
                    return row[0]

            cursor.execute(
                """
                select id
                from curriculum."Training_plan"
                where lower(%s) like ('%%' || lower(module_name) || '%%')
                   or lower(%s) like ('%%' || lower(module_name) || '%%')
                order by length(module_name) desc, id desc
                limit 1
                """,
                [module_text, title_text],
            )
            row = cursor.fetchone()
            return row[0] if row else None
    except Exception:
        return None


def _match_programme_catalogue_id(programme, module, title, supplied_id=None):
    supplied_text = str(supplied_id or "").strip()
    if supplied_text and not _is_int_like(supplied_text):
        return supplied_text

    candidates = _program_candidates(programme)
    module_text = module or ""
    title_text = title or ""

    try:
        with connection.cursor() as cursor:
            if supplied_text:
                cursor.execute(
                    """
                    select programme_id
                    from curriculum.module_authoring_modules
                    where coalesce(trim(programme_id), '') <> ''
                      and imported_from_training_plan_id = %s
                    order by updated_at desc nulls last, created_at desc nulls last
                    limit 1
                    """,
                    [supplied_text],
                )
                row = cursor.fetchone()
                if row and row[0]:
                    return row[0]

                cursor.execute(
                    """
                    select mam.programme_id
                    from curriculum."Training_plan" tp
                    join curriculum.module_authoring_modules mam
                      on coalesce(trim(mam.programme_id), '') <> ''
                     and (
                       lower(mam.programme_name) = lower(tp."Program")
                       or lower(mam.programme_id) = lower(tp."Program")
                     )
                    where tp.id = %s
                    order by
                      case
                        when lower(mam.title) = lower(%s) then 0
                        when lower(%s) like ('%%' || lower(mam.title) || '%%') then 1
                        else 2
                      end,
                      mam.updated_at desc nulls last,
                      mam.created_at desc nulls last
                    limit 1
                    """,
                    [supplied_text, module_text, module_text],
                )
                row = cursor.fetchone()
                if row and row[0]:
                    return row[0]

            programme_clauses = []
            params = []
            for candidate in candidates:
                programme_clauses.append(
                    "(lower(programme_name) = lower(%s) or lower(programme_id) = lower(%s) or lower(%s) like lower(programme_name) || ' %%')"
                )
                params.extend([candidate, candidate, candidate])

            where_programme = f"and ({' or '.join(programme_clauses)})" if programme_clauses else ""
            cursor.execute(
                f"""
                select programme_id
                from curriculum.module_authoring_modules
                where coalesce(trim(programme_id), '') <> ''
                  {where_programme}
                order by
                  case
                    when lower(title) = lower(%s) then 0
                    when lower(%s) like ('%%' || lower(title) || '%%') then 1
                    when lower(%s) like ('%%' || lower(title) || '%%') then 2
                    else 3
                  end,
                  length(title) desc,
                  updated_at desc nulls last,
                  created_at desc nulls last
                limit 1
                """,
                [*params, module_text, module_text, title_text],
            )
            row = cursor.fetchone()
            return row[0] if row else (supplied_text or "")
    except Exception:
        return supplied_text or ""


def _training_plan_programme_for_id(plan_id):
    if not plan_id or not _is_int_like(plan_id):
        return ""
    try:
        with connection.cursor() as cursor:
            cursor.execute(
                'select "Program" from curriculum."Training_plan" where id = %s limit 1',
                [plan_id],
            )
            row = cursor.fetchone()
            return row[0] if row else ""
    except Exception:
        return ""


def _training_plan_programmes():
    try:
        with connection.cursor() as cursor:
            cursor.execute(
                """
                select "Program", count(*) as rows_count
                from curriculum."Training_plan"
                where coalesce(trim("Program"), '') <> ''
                group by "Program"
                order by "Program"
                """
            )
            return [{"name": row[0], "trainingPlanRows": row[1]} for row in cursor.fetchall()]
    except Exception:
        return []


def _training_plan_programme_map(plan_ids):
    ids = [int(plan_id) for plan_id in plan_ids if _is_int_like(plan_id)]
    if not ids:
        return {}
    try:
        with connection.cursor() as cursor:
            placeholders = ", ".join(["%s"] * len(ids))
            cursor.execute(
                f'select id, "Program" from curriculum."Training_plan" where id in ({placeholders})',
                ids,
            )
            return {row[0]: row[1] for row in cursor.fetchall()}
    except Exception:
        return {}


def _ensure_quiz_course_links_table():
    with connection.cursor() as cursor:
        cursor.execute(
            """
            create table if not exists curriculum.quiz_course_links (
              id bigserial primary key,
              quiz_id bigint not null references curriculum.quizzes(id) on delete cascade,
              training_plan_id integer not null references curriculum."Training_plan"(id) on delete cascade,
              created_at timestamptz not null default now(),
              unique (quiz_id, training_plan_id)
            )
            """
        )


def _quiz_course_link_ids(quiz_id):
    _ensure_quiz_course_links_table()
    with connection.cursor() as cursor:
        cursor.execute(
            "select training_plan_id from curriculum.quiz_course_links where quiz_id = %s order by training_plan_id",
            [quiz_id],
        )
        return [row[0] for row in cursor.fetchall()]


def _training_plan_courses_for_programme(programme):
    if not programme:
        return []
    with connection.cursor() as cursor:
        cursor.execute(
            """
            select
              id,
              "Program",
              module_name,
              "Cohort_name",
              "Starting_date_lable"
            from curriculum."Training_plan"
            where lower(trim("Program")) = lower(trim(%s))
              and coalesce(trim(module_name), '') <> ''
            order by module_name, "Cohort_name", id
            """,
            [programme],
        )
        rows = cursor.fetchall()

    courses = []
    seen = set()
    for plan_id, program, module_name, cohort_name, starting_date in rows:
        if _is_placeholder_training_value(program) or _is_placeholder_training_value(module_name):
            continue
        key = (module_name, cohort_name or "", starting_date or "")
        if key in seen:
            continue
        seen.add(key)
        label_parts = [module_name]
        if cohort_name:
            label_parts.append(cohort_name)
        if starting_date:
            label_parts.append(str(starting_date))
        courses.append({
            "id": plan_id,
            "programme": program,
            "module": module_name,
            "cohort": cohort_name or "",
            "startDate": str(starting_date or ""),
            "label": " - ".join(label_parts),
        })
    return courses


def _is_placeholder_training_value(value):
    text = (value or "").lower()
    return "test" in text or "delete" in text


@require_http_methods(["GET"])
def training_plan_options(request):
    try:
        with connection.cursor() as cursor:
            cursor.execute(
                """
                select
                  tp."Program",
                  tp.module_name,
                  max(coalesce(nullif(mam.programme_id, ''), tp."Program")) as programme_id,
                  max(tp.id) as training_plan_id,
                  max(mam.module_catalogue_id) as module_catalogue_id
                from curriculum."Training_plan" tp
                left join curriculum.module_authoring_modules mam
                  on mam.imported_from_training_plan_id = tp.id::text
                  or (
                    mam.title = tp.module_name
                    and (
                      mam.programme_id = tp."Program"
                      or mam.programme_name = tp."Program"
                      or tp."Program" like mam.programme_name || ' %%'
                    )
                  )
                where coalesce(trim(tp."Program"), '') <> ''
                  and coalesce(trim(tp.module_name), '') <> ''
                group by tp."Program", tp.module_name
                order by tp."Program", tp.module_name
                """
            )
            rows = cursor.fetchall()
    except Exception:
        rows = []

    programmes = []
    modules_by_programme = {}
    seen_programmes = set()
    for programme, module_name, programme_id, training_plan_id, module_catalogue_id in rows:
        if _is_placeholder_training_value(programme) or _is_placeholder_training_value(module_name):
            continue
        if programme not in seen_programmes:
            programmes.append({"value": programme, "label": programme})
            seen_programmes.add(programme)
        modules_by_programme.setdefault(programme, []).append({
            "value": module_name,
            "label": module_name,
            "programmeId": programme_id,
            "trainingPlanId": training_plan_id,
            "moduleId": module_catalogue_id or "",
        })

    return JsonResponse({
        "programmes": programmes,
        "modulesByProgramme": modules_by_programme,
    })


def _package_type_from_extension(extension):
    mapping = {
        ".xml": "xml",
        ".zip": "scorm",
        ".scorm": "scorm",
        ".xlsx": "excel",
        ".xlsm": "excel",
        ".csv": "csv",
    }
    return mapping.get(extension, "file")


def _normalise_correct_answer(value):
    if value is None:
        return ""
    text = str(value).strip()
    match = re.match(r"^([A-Z])(?:[\).:\s]|$)", text, flags=re.IGNORECASE)
    return match.group(1).upper() if match else text.lower()


def _correct_answer_index(value):
    key = _normalise_correct_answer(value)
    if key.isdigit():
        return max(int(key) - 1, 0)
    if len(key) == 1 and key.isalpha():
        return ord(key.upper()) - ord("A")
    return 0


def _question_from_parts(text, options, correct_answer="", explanation="", question_type="single_choice"):
    question_type = _normalise_question_type(question_type or "single_choice")
    cleaned_options = [str(option).strip() for option in options if option is not None and str(option).strip()]
    correct_tokens = [
        _normalise_correct_answer(token)
        for token in re.split(r"[,;/|]", str(correct_answer or ""))
        if str(token).strip()
    ]
    correct_keys = set(correct_tokens)
    correct_key = next(iter(correct_keys), "")
    answers = []
    for index, option in enumerate(cleaned_options):
        letter = chr(65 + index)
        number = str(index + 1)
        option_key = _normalise_correct_answer(option)
        answers.append({
            "text": re.sub(r"^[A-Z][\).:\s]+", "", option, flags=re.IGNORECASE).strip(),
            "is_correct": bool({letter, number, option_key, option.lower()} & correct_keys),
            "sort_order": index,
        })

    if answers and correct_key and not any(answer["is_correct"] for answer in answers):
        answers[0]["is_correct"] = True

    return {
        "text": text.strip(),
        "question_type": question_type,
        "answers": answers,
        "explanation": explanation.strip(),
    }


def _parse_questions_from_text(text):
    questions = []
    current_text = ""
    current_options = []
    current_answer = ""
    current_explanation = ""

    def flush():
        nonlocal current_text, current_options, current_answer, current_explanation
        if current_text:
            questions.append(_question_from_parts(current_text, current_options, current_answer, current_explanation))
        current_text = ""
        current_options = []
        current_answer = ""
        current_explanation = ""

    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue

        question_match = re.match(r"^(?:Q(?:uestion)?\s*)?\d+[\).:-]\s*(.+)$", line, flags=re.IGNORECASE)
        if question_match:
            flush()
            current_text = question_match.group(1).strip()
            continue

        q_prefix = re.match(r"^Q(?:uestion)?[\s:.-]+(.+)$", line, flags=re.IGNORECASE)
        if q_prefix:
            flush()
            current_text = q_prefix.group(1).strip()
            continue

        option_match = re.match(r"^([A-H])[\).:-]\s*(.+)$", line, flags=re.IGNORECASE)
        if option_match:
            current_options.append(f"{option_match.group(1).upper()}. {option_match.group(2).strip()}")
            continue

        answer_match = re.match(r"^(?:answer|correct\s*answer|correct)[:\s-]+(.+)$", line, flags=re.IGNORECASE)
        if answer_match:
            current_answer = answer_match.group(1).strip()
            continue

        explanation_match = re.match(r"^(?:explanation|feedback)[:\s-]+(.+)$", line, flags=re.IGNORECASE)
        if explanation_match:
            current_explanation = explanation_match.group(1).strip()
            continue

        if current_options:
            current_options[-1] = f"{current_options[-1]} {line}"
        elif current_text:
            current_text = f"{current_text} {line}"

    flush()
    return [question for question in questions if question["text"]]


def _parse_questions_from_xml(root):
    parsed = []
    record_nodes = root.findall(".//record")
    if record_nodes:
        for record_node in record_nodes:
            question_text = record_node.findtext("Question_Title") or record_node.findtext("question_title") or ""
            if not question_text:
                continue
            correct_answer = record_node.findtext("Answer") or record_node.findtext("answer") or ""
            explanation_index = _correct_answer_index(correct_answer)
            explanation_tag = "Question_Explanation" if explanation_index == 0 else f"Question_Explanation_{explanation_index}"
            options = [
                record_node.findtext("Option_1"),
                record_node.findtext("Option_2"),
                record_node.findtext("Option_3"),
                record_node.findtext("Option_4"),
                record_node.findtext("Option_5"),
            ]
            parsed.append(_question_from_parts(
                question_text,
                options,
                correct_answer,
                record_node.findtext(explanation_tag) or "",
                record_node.findtext("Question_Type") or "single_choice",
            ))
        return parsed

    for question_node in root.findall(".//question"):
        text = (
            question_node.findtext("stem")
            or question_node.findtext("text")
            or question_node.findtext("question_text")
            or question_node.get("text")
            or ""
        )
        options = []
        correct_answer = ""
        for index, option_node in enumerate(question_node.findall(".//option")):
            option_text = "".join(option_node.itertext()).strip()
            options.append(option_text)
            if option_node.get("correct", "").lower() == "true":
                correct_answer = chr(65 + index)
        if text:
            parsed.append(_question_from_parts(text, options, correct_answer, question_node.findtext("feedback") or ""))
    return parsed


def _normalise_question_type(value):
    text = str(value or "single_choice").strip().lower().replace("-", "_").replace(" ", "_")
    aliases = {
        "single": "single_choice",
        "singlechoice": "single_choice",
        "multi_choice": "multiple_choice",
        "multiple": "multiple_choice",
        "multiplechoice": "multiple_choice",
        "truefalse": "true_false",
        "true_false": "true_false",
        "fill_blank": "fill_gap",
        "fill_in_the_gap": "fill_gap",
        "fillinthegap": "fill_gap",
        "image": "image_matching",
        "imagematching": "image_matching",
        "match": "matching",
        "keyword": "keywords",
        "order": "ordering",
    }
    return aliases.get(text, text or "single_choice")


def _parse_questions_from_course_data(data):
    parsed = []
    questions = data.get("questions", []) if isinstance(data, dict) else []
    if not isinstance(questions, list):
        return parsed

    for question in questions:
        if not isinstance(question, dict):
            continue
        text = question.get("title") or question.get("text") or question.get("question") or ""
        options = question.get("options", [])
        option_texts = []
        correct_answer = question.get("answer") or question.get("correctAnswer") or question.get("correct_answer") or ""
        explanation = question.get("explanation") or question.get("feedback") or ""

        if isinstance(options, list):
            for index, option in enumerate(options):
                if isinstance(option, dict):
                    option_texts.append(option.get("text") or option.get("title") or option.get("label") or "")
                    option_id = str(option.get("id", "")).strip()
                    if not explanation and correct_answer and option_id == str(correct_answer).strip():
                        explanation = option.get("explanation") or ""
                    if not correct_answer and option.get("correct") is True:
                        correct_answer = option_id or str(index + 1)
                        explanation = option.get("explanation") or explanation
                else:
                    option_texts.append(str(option))

        if text:
            parsed.append(_question_from_parts(
                str(text),
                option_texts,
                correct_answer,
                str(explanation),
                _normalise_question_type(question.get("type") or question.get("questionType") or question.get("question_type")),
            ))
    return parsed


def _extract_json_object_from_javascript(text):
    object_start = text.find("{")
    if object_start == -1:
        return None
    decoder = json.JSONDecoder()
    try:
        data, _ = decoder.raw_decode(text[object_start:])
        return data
    except json.JSONDecodeError:
        return None


def _extract_excel_questions(uploaded_file):
    from openpyxl import load_workbook

    uploaded_file.seek(0)
    workbook = load_workbook(uploaded_file, read_only=True, data_only=True)
    sheet = workbook.active
    rows = list(sheet.iter_rows(values_only=True))
    if not rows:
        return []

    headers = [str(value).strip().lower() if value is not None else "" for value in rows[0]]
    header_map = {header: index for index, header in enumerate(headers)}

    def get(row, *names):
        for name in names:
            if name in header_map and header_map[name] < len(row):
                return row[header_map[name]]
        return None

    parsed = []
    detected_title = ""
    metadata = {}
    for row in rows[1:]:
        question_text = get(row, "question", "question_title", "question title", "question_text", "question text", "stem")
        if not question_text:
            continue
        detected_title = detected_title or str(get(row, "quiz_title", "quiz title") or "").strip()
        metadata.setdefault("programme", str(get(row, "course_name", "course name") or "").strip())
        metadata.setdefault("module", str(get(row, "section_name", "section name") or "").strip())
        options = [
            get(row, "option_1", "option 1", "option_a", "option a", "a"),
            get(row, "option_2", "option 2", "option_b", "option b", "b"),
            get(row, "option_3", "option 3", "option_c", "option c", "c"),
            get(row, "option_4", "option 4", "option_d", "option d", "d"),
            get(row, "option_5", "option 5", "option_e", "option e"),
        ]
        correct_answer = get(row, "correct_answer", "correct answer", "answer", "correct")
        explanation_index = _correct_answer_index(correct_answer)
        explanation_names = ["question explanation"]
        if explanation_index > 0:
            explanation_names = [f"question explanation_{explanation_index}", f"question_explanation_{explanation_index}"]
        explanation = get(row, *explanation_names, "explanation", "feedback") or ""
        question_type = str(get(row, "question_type", "question type") or "single_choice").strip()
        parsed.append(_question_from_parts(str(question_text), options, correct_answer, str(explanation), question_type))
    return parsed, detected_title, metadata


def _extract_csv_questions(uploaded_file):
    text = _decode_uploaded_text(uploaded_file)
    reader = csv.DictReader(text.splitlines())
    rows = list(reader)
    parsed = []
    detected_title = ""
    metadata = {}

    field_names = {str(field or "").strip().lower() for field in (reader.fieldnames or [])}
    grouped_answer_format = {"question_id", "question_text", "option_text", "is_correct"}.issubset(field_names)
    if grouped_answer_format:
        grouped = {}
        for row in rows:
            normalised = {key.strip().lower(): value for key, value in row.items() if key}
            question_id = str(normalised.get("question_id") or normalised.get("id") or "").strip()
            question_text = _clean_csv_text(normalised.get("question_text") or normalised.get("question") or "")
            option_text = _clean_csv_text(normalised.get("option_text") or normalised.get("answer_text") or "")
            if not question_id or not question_text or not option_text:
                continue
            item = grouped.setdefault(question_id, {
                "text": question_text,
                "type": _normalise_question_type(normalised.get("question_type") or "single_choice"),
                "options": [],
                "correct_letters": [],
                "explanation": _clean_csv_text(normalised.get("explanation") or normalised.get("feedback") or ""),
                "title": str(normalised.get("quiz_title") or normalised.get("title") or "").strip(),
            })
            answer_index = len(item["options"])
            item["options"].append(option_text)
            if _csv_truthy(normalised.get("is_correct")):
                item["correct_letters"].append(chr(65 + answer_index))
            if not item["explanation"]:
                item["explanation"] = _clean_csv_text(normalised.get("explanation") or normalised.get("feedback") or "")

        for item in grouped.values():
            detected_title = detected_title or item["title"]
            correct_answer = ", ".join(item["correct_letters"])
            parsed.append(_question_from_parts(item["text"], item["options"], correct_answer, item["explanation"], item["type"]))
        uploaded_file.seek(0)
        return parsed, detected_title, metadata

    for row in rows:
        normalised = {key.strip().lower(): value for key, value in row.items() if key}

        def get(*names):
            for name in names:
                key = name.strip().lower()
                variants = {
                    key,
                    key.replace("_", " "),
                    key.replace(" ", "_"),
                    key.replace("-", "_"),
                }
                for variant in variants:
                    if variant in normalised and str(normalised[variant]).strip():
                        return normalised[variant]
            return ""

        question_text = _clean_csv_text(get("question title", "question", "question_text", "question text", "stem"))
        if not question_text:
            continue
        detected_title = detected_title or str(get("quiz_title", "quiz title", "title")).strip()
        metadata.setdefault("programme", str(get("programme", "course_name", "course name")).strip())
        metadata.setdefault("module", str(get("module", "section_name", "section name", "part")).strip())
        question_type = _normalise_question_type(get("question_type", "question type") or "single_choice")
        options = [
            get("option 1", "option_1", "option_a", "option a", "a"),
            get("option 2", "option_2", "option_b", "option b", "b"),
            get("option 3", "option_3", "option_c", "option c", "c"),
            get("option 4", "option_4", "option_d", "option d", "d"),
            get("option 5", "option_5", "option_e", "option e"),
        ]
        correct_answer = get("correct_answer", "correct answer", "answer", "correct")
        if question_type in {"matching", "image_matching"}:
            pairs = get("matching_pairs", "matching pairs")
            options = [item.strip().replace(" = ", " -> ") for item in re.split(r"\s*\|\s*", pairs) if item.strip()]
            correct_answer = correct_answer or ", ".join(chr(65 + index) for index in range(len(options)))
        elif question_type == "keywords":
            keywords = get("keywords")
            options = [item.strip() for item in re.split(r"\s*[;|]\s*", keywords) if item.strip()]
            correct_answer = correct_answer or ", ".join(chr(65 + index) for index in range(len(options)))
        elif question_type == "ordering":
            ordering_items = get("ordering_items", "ordering items")
            options = [item.strip() for item in re.split(r"\s*(?:>|;|\|)\s*", ordering_items) if item.strip()]
            correct_answer = correct_answer or ", ".join(chr(65 + index) for index in range(len(options)))
        elif question_type == "fill_gap" and not any(str(option or "").strip() for option in options) and correct_answer:
            options = [correct_answer]
        explanation_index = _correct_answer_index(correct_answer)
        explanation_key = "question explanation" if explanation_index == 0 else f"question explanation_{explanation_index}"
        explanation = _clean_csv_text(get(explanation_key, "question explanation", "explanation", "feedback"))
        parsed.append(_question_from_parts(question_text, options, correct_answer, explanation, question_type))
    uploaded_file.seek(0)
    return parsed, detected_title, metadata


def _validate_xml(uploaded_file):
    uploaded_file.seek(0)
    try:
        tree = ElementTree.parse(uploaded_file)
        root = tree.getroot()
        parsed_questions = _parse_questions_from_xml(root)
        questions = len(parsed_questions) or len(root.findall(".//question"))
        title_node = root.find(".//title")
        record_node = root.find(".//record")
        title = title_node.text.strip() if title_node is not None and title_node.text else ""
        if not title and record_node is not None:
            title = record_node.findtext("Quiz_Title") or ""
        metadata = {}
        if record_node is not None:
            metadata = {
                "programme": record_node.findtext("Course_Name") or "",
                "module": record_node.findtext("Section_Name") or "",
            }
        return True, "", questions, title.strip(), parsed_questions, metadata
    except ElementTree.ParseError as exc:
        return False, f"XML parse error: {exc}", 0, "", [], {}
    finally:
        uploaded_file.seek(0)


def _validate_scorm(uploaded_file):
    uploaded_file.seek(0)
    try:
        with zipfile.ZipFile(uploaded_file) as package:
            names = package.namelist()
            has_manifest = any(Path(name).name.lower() == "imsmanifest.xml" for name in names)
            message = "" if has_manifest else "SCORM package is missing imsmanifest.xml"
            parsed_questions = []
            detected_title = ""
            metadata = {}

            for name in names:
                lower_name = name.lower()
                if lower_name.endswith("/") or lower_name.endswith((".png", ".jpg", ".jpeg", ".gif", ".css")):
                    continue

                try:
                    raw = package.read(name)
                    text = raw.decode("utf-8-sig", errors="replace")
                except Exception:
                    continue

                if lower_name.endswith(".xml") and Path(name).name.lower() != "imsmanifest.xml":
                    try:
                        root = ElementTree.fromstring(text)
                    except ElementTree.ParseError:
                        root = None
                    if root is not None:
                        parsed_questions = _parse_questions_from_xml(root)
                        if parsed_questions:
                            title_node = root.find(".//title")
                            detected_title = title_node.text.strip() if title_node is not None and title_node.text else detected_title
                            break

                if lower_name.endswith((".js", ".json")):
                    try:
                        data = json.loads(text) if lower_name.endswith(".json") else _extract_json_object_from_javascript(text)
                    except json.JSONDecodeError:
                        data = None
                    if isinstance(data, dict):
                        questions = _parse_questions_from_course_data(data)
                        if questions:
                            parsed_questions = questions
                            detected_title = data.get("quizTitle") or data.get("title") or detected_title
                            metadata = {
                                "programme": data.get("courseName") or "",
                                "module": data.get("sectionName") or "",
                            }
                            break

                if lower_name.endswith(".txt") and not parsed_questions:
                    parsed_questions = _parse_questions_from_text(text)
                    if parsed_questions:
                        break

            question_count = len(parsed_questions)
            if has_manifest and not parsed_questions:
                message = "SCORM package uploaded, but no readable quiz questions were found inside it."
            return has_manifest, message, question_count, str(detected_title).strip(), parsed_questions, metadata
    except zipfile.BadZipFile:
        return False, "SCORM upload must be a valid .zip package", 0, "", [], {}
    finally:
        uploaded_file.seek(0)


def _manifest_launch_path(package):
    for name in package.namelist():
        if Path(name).name.lower() != "imsmanifest.xml":
            continue
        try:
            root = ElementTree.fromstring(package.read(name))
        except ElementTree.ParseError:
            return ""
        for resource in root.findall(".//{*}resource"):
            href = resource.get("href", "").strip()
            if href:
                return href
    return ""


def _safe_scorm_member_path(destination, member_name):
    cleaned = normpath(member_name.replace("\\", "/")).lstrip("/")
    if not cleaned or cleaned == "." or cleaned.startswith("../"):
        return None
    target = (destination / cleaned).resolve()
    resolved_destination = destination.resolve()
    if target != resolved_destination and resolved_destination not in target.parents:
        return None
    return target


def _ensure_scorm_extracted(quiz):
    if quiz.package_type != "scorm" or not _scorm_uploaded_file_readable(quiz):
        raise Http404

    destination = (settings.MEDIA_ROOT / "scorm_runtime" / f"quiz_{quiz.id}").resolve()
    marker = destination / ".source"
    launch_marker = destination / ".launch"
    source_key = f"{quiz.file_name}:{quiz.file_size}:{quiz.updated_at.isoformat()}"

    if marker.exists() and marker.read_text(encoding="utf-8") == source_key:
        launch_path = launch_marker.read_text(encoding="utf-8") if launch_marker.exists() else "index.html"
        return destination, launch_path or "index.html"

    destination.mkdir(parents=True, exist_ok=True)
    for existing in sorted(destination.rglob("*"), key=lambda path: len(path.parts), reverse=True):
        if existing.is_file():
            existing.unlink()
        elif existing.is_dir() and existing != destination:
            try:
                existing.rmdir()
            except OSError:
                pass

    try:
        with quiz.uploaded_file.open("rb") as uploaded:
            with zipfile.ZipFile(uploaded) as package:
                launch_path = _manifest_launch_path(package)
                for member in package.infolist():
                    target = _safe_scorm_member_path(destination, member.filename)
                    if target is None:
                        continue
                    if member.is_dir():
                        target.mkdir(parents=True, exist_ok=True)
                        continue
                    target.parent.mkdir(parents=True, exist_ok=True)
                    with package.open(member) as source, target.open("wb") as output:
                        output.write(source.read())

                if not launch_path:
                    html_files = [name for name in package.namelist() if name.lower().endswith((".html", ".htm"))]
                    index_file = next((name for name in html_files if Path(name).name.lower() == "index.html"), "")
                    launch_path = index_file or (html_files[0] if html_files else "")
    except zipfile.BadZipFile:
        raise Http404

    marker.write_text(source_key, encoding="utf-8")
    launch_marker.write_text(launch_path or "index.html", encoding="utf-8")
    return destination, launch_path or "index.html"


def _extract_questions_from_upload(uploaded_file, package_type):
    try:
        if package_type == "excel":
            questions, detected_title, metadata = _extract_excel_questions(uploaded_file)
        elif package_type == "csv":
            questions, detected_title, metadata = _extract_csv_questions(uploaded_file)
        else:
            return True, "", 0, "", [], {}
    except Exception as exc:
        uploaded_file.seek(0)
        return False, f"Could not read questions from file: {exc}", 0, "", [], {}

    uploaded_file.seek(0)
    return True, "", len(questions), detected_title, questions, metadata


def _decode_uploaded_text(uploaded_file):
    uploaded_file.seek(0)
    raw = uploaded_file.read()
    uploaded_file.seek(0)
    last_error = None
    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError as exc:
            last_error = exc
    if last_error:
        raise last_error
    return raw.decode("utf-8-sig", errors="replace")


def _clean_csv_text(value):
    text = html.unescape(str(value or "").replace("\xa0", " "))
    text = re.sub(r"<br\s*/?>", "\n", text, flags=re.IGNORECASE)
    text = re.sub(r"</p\s*>", "\n", text, flags=re.IGNORECASE)
    text = re.sub(r"<[^>]+>", "", text)
    return re.sub(r"[ \t\r\f\v]+", " ", text).strip()


def _csv_truthy(value):
    return str(value or "").strip().lower() in {"yes", "true", "1", "y", "correct"}


def _save_questions(quiz, questions):
    for index, question in enumerate(questions):
        question_record = QuizQuestion.objects.create(
            quiz=quiz,
            question_text=question["text"],
            question_type=question.get("question_type", quiz.default_question_type),
            explanation=question.get("explanation", ""),
            sort_order=index,
        )
        for answer in question.get("answers", []):
            QuizAnswer.objects.create(
                question=question_record,
                answer_text=answer["text"],
                is_correct=answer["is_correct"],
                sort_order=answer["sort_order"],
            )


def _default_question_type_from_questions(questions, fallback="single_choice"):
    first_type = next((question.get("question_type") for question in questions if question.get("question_type")), "")
    return _normalise_question_type(first_type or fallback)


def _seed_quizzes():
    if QuizPackage.objects.exists():
        return

    records = [
        ("Business Communication - Week 1 Quiz", "Business Communication", "Business Admin L3", "v2.1", 12, "published", "xml", 24576, True, 2, "Jack Williams", 1),
        ("Written Communication Assessment", "Business Communication", "Business Admin L3", "v1.8", 15, "published", "xml", 31744, True, 3, "Jack Williams", 1),
        ("Organisational Culture Checkpoint", "Organisational Culture", "Business Admin L3", "v1.5", 10, "draft", "xml", 18432, True, 1, "Fatma Mo", 1),
        ("Data Visualisation - Tableau Basics", "Data Visualisation", "Data Analyst L4", "v2.0", 18, "published", "xml", 43008, True, 2, "Rewan yasser", 2),
        ("Statistical Concepts Quiz", "Statistical Analysis", "Data Analyst L4", "v1.3", 20, "validating", "xml", 46080, False, 3, "Jack Williams", 2),
        ("Segmentation & Targeting Test", "Marketing Planning", "Marketing Exec L4", "v1.7", 14, "published", "scorm", 28672, True, 2, "Fatma Mo", 1),
        ("Digital Channels Assessment", "Digital Channels", "Marketing Exec L4", "v1.4", 16, "draft", "xml", 33792, True, 2, "Jack Williams", 1),
        ("Agile Development Fundamentals", "Agile Development", "Software Dev L4", "v0.9", 22, "draft", "scorm", 51200, False, 4, "Jack Williams", 1),
    ]

    now = timezone.now()
    for title, module, programme, version, questions, status, package_type, file_size, schema_valid, mapped, author, courses in records:
        QuizPackage.objects.create(
            title=title,
            module=module,
            programme=programme,
            version=version,
            questions=questions,
            status=status,
            package_type=package_type,
            file_name=f"{title.lower().replace(' ', '-')}.{package_type if package_type == 'xml' else 'zip'}",
            file_size=file_size,
            schema_valid=schema_valid,
            validation_message="" if schema_valid else "Package requires schema review before publishing.",
            mapped_components=mapped,
            author=author,
            linked_courses=courses,
            published_at=now if status == "published" else None,
        )


@csrf_exempt
@require_http_methods(["POST"])
def generate_ai_questions(request):
    if not settings.OPENAI_API_KEY:
        return JsonResponse({"error": "OPENAI_API_KEY is not configured in backend .env"}, status=503)

    readable_files = []
    unreadable_files = []
    try:
        if request.FILES:
            uploaded_files = request.FILES.getlist("files") or request.FILES.getlist("file")
            source_text, readable_files, unreadable_files = _extract_text_from_ai_files_with_report(uploaded_files)
            payload = request.POST
            fallback_text = str(payload.get("lessonContent") or payload.get("lesson_content") or payload.get("text") or "").strip()
            source_text = source_text.strip() or fallback_text
        else:
            body = json.loads(request.body or "{}")
            source_text = str(body.get("lessonContent") or body.get("lesson_content") or body.get("text") or "").strip()
            payload = body
    except json.JSONDecodeError:
        return JsonResponse({"error": "Invalid JSON body"}, status=400)
    except ValueError as exc:
        return JsonResponse({"error": str(exc)}, status=400)

    topic = str(payload.get("topic") or "").strip()
    custom_instructions = str(
        payload.get("customInstructions")
        or payload.get("custom_instructions")
        or payload.get("customPrompt")
        or payload.get("custom_prompt")
        or ""
    ).strip()
    if not source_text and not topic:
        unreadable_note = f" Unreadable files: {', '.join(unreadable_files)}." if unreadable_files else ""
        return JsonResponse({
            "error": f"No readable lesson content found. Add a topic, paste lesson text, or upload text-based PDF/PowerPoint/Word files.{unreadable_note} The embedded prompt is instructions only, not lesson content.",
            "unreadableFiles": unreadable_files,
        }, status=400)

    try:
        question_count = int(payload.get("questionCount") or payload.get("question_count") or 5)
    except (TypeError, ValueError):
        question_count = 5
    question_count = max(1, min(question_count, 60))

    course_title = str(payload.get("courseTitle") or payload.get("course_title") or "").strip()
    module = str(payload.get("module") or "").strip()
    programme = str(payload.get("programme") or "").strip()
    source_text = source_text[:AI_GENERATION_MAX_CONTENT_CHARS]
    source_files = _source_file_names_from_text(source_text)
    source_file_count = len(source_files)
    source_allocation = _question_allocation_instruction(source_files, question_count)

    schema = {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "questions": {
                "type": "array",
                "minItems": 1,
                "maxItems": question_count,
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "question": {"type": "string"},
                        "type": {"type": "string", "enum": sorted(AI_GENERATION_QUESTION_TYPES)},
                        "options": {
                            "type": "array",
                            "minItems": 2,
                            "maxItems": 8,
                            "items": {"type": "string"},
                        },
                        "correct_answer": {
                            "type": "string",
                            "description": "For single_choice/true_false/fill_gap use one option letter or exact option text. For multiple_choice use 2 or 3 correct option letters separated by commas. For matching/image_matching/keywords/ordering, options must contain the correct items and correct_answer may be ALL or the item letters.",
                        },
                        "explanation": {"type": "string"},
                        "difficulty": {"type": "string", "enum": ["easy", "medium", "hard"]},
                    },
                    "required": ["question", "type", "options", "correct_answer", "explanation", "difficulty"],
                },
            }
        },
        "required": ["questions"],
    }

    prompt = f"""
Act like an assessment design SME in adult learning and an instructional designer creating professional, high-quality mixed-format LMS quizzes aligned to KSBs (Knowledge, Skills, Behaviours).

Generate high-quality LMS quiz questions for a curriculum admin to review before saving.

Course title: {course_title or "Not provided"}
Programme: {programme or "Not provided"}
Module: {module or "Not provided"}
Topic: {topic or "Use the source content"}
Difficulty: progressive mix of easy, medium and hard
Question types: mixed. Use a suitable spread across single_choice, multiple_choice, true_false, matching, image_matching, keywords, fill_gap and ordering when the source supports them.
Number of questions: {question_count}
Detected source files: {source_file_count or "Not provided"}
{source_allocation}

Source content:
{source_text or topic}

Editable embedded prompt from admin UI:
{custom_instructions or "None provided. Follow the embedded core prompt only."}

Rules:
- Return exactly {question_count} questions when the source content supports it.
- Use ALL supplied source files together. Do not base the quiz on only one file, only the last file, or only the longest/easiest file.
- Treat each "Source file:" block as a separate required content source.
- Follow the required source-file allocation above exactly when source-file names are provided.
- If the allocation says a file needs 3 questions, exactly 3 generated questions must be grounded in that file's concepts.
- Before writing questions, mentally identify the main concepts in every source file block.
- Coverage rule: every source file that contains readable lesson content must be represented by at least one question when {question_count} is greater than or equal to the number of readable source files.
- If {question_count} is smaller than the number of readable source files, cover as many different source files as possible and prioritise the most assessment-worthy concept from each.
- Distribute questions as evenly as possible across source files before adding extra questions to richer files.
- For 10 questions across 3 readable files, aim for a 3/3/4 or 3/4/3 split unless the content clearly requires otherwise.
- For 20 questions across 3 readable files, aim for roughly 6/7/7.
- For 60 questions across 3 readable files, aim for 20 questions per file/section.
- Do not generate all questions from Part 2 or Part 3 if Part 1 is present; each part must contribute.
- If any uploaded XML/CSV looks like an existing quiz, treat it as a style/format example unless the user explicitly asks to copy it.
- Use ONLY the supplied source content and uploaded examples. Do not use external knowledge.
- Write questions for adult UK workplace LMS learners aged around 30-50 at the programme level shown above.
- If the user uploads an example exam/quiz XML, match its level, tone, structure, and option style while generating new questions from the learning materials.
- Identify the main KSBs from the source content and mentally map each question to one or more K/S/B items.
- Mention the targeted KSB(s) briefly inside each explanation, for example: "KSBs covered: K2.1, S3.2" when source KSB labels are available.
- If KSB labels are not explicit, use concise inferred labels such as "Knowledge: AI governance" or "Skill: marketing judgement" in the explanation.
- Prefer applied understanding and judgement questions over simple recall.
- At least half of the questions should use a short realistic workplace/marketing scenario.
- Cover distinct concepts across the source material; avoid multiple questions testing the same idea.
- When {question_count} is 60, organise coverage internally into 3 balanced themes/sections of 20 questions each, based on the source content.
- When {question_count} is 60, balance correct answer letters across the full quiz as evenly as possible: 15 A, 15 B, 15 C, 15 D.
- When {question_count} is 20, balance correct answer letters as 5 A, 5 B, 5 C, 5 D.
- Avoid obvious answer-letter patterns and avoid more than two identical correct letters in a row.
- Use a mixed set of question types. Do not generate the whole quiz as only single_choice unless the requested count is 1 and the source cannot support another type.
- Include a progressive difficulty curve across the quiz: start with easier recall/comprehension questions, move into medium application questions, and include harder judgement/scenario questions near the end. Label each question difficulty as easy, medium or hard.
- For counts of 8 or more, include at least 4 different question types where the source supports them.
- For counts of 12 or more, include at least 6 different question types where the source supports them.
- For single_choice questions, provide 4 plausible options and exactly one correct answer.
- For multiple_choice questions, provide 4 plausible options and exactly 2 or 3 correct answers. The correct_answer field must list all correct option letters or exact option texts separated by commas, for example "A, C" or "A, B, D". Do not create multiple_choice questions with only one correct answer.
- For matching questions, write the stem as a matching task and put every correct pair in options using "left -> right" format, for example "Instagram -> visual discovery". Do not put the pairs only in explanation.
- For image_matching questions, write a visual/image-selection style stem and put every image/concept match in options using "image placeholder -> concept" format if no actual image assets exist.
- For keywords questions, ask learners to provide or identify key terms; put every accepted keyword or short concept as a separate option. Do not put accepted keywords only in explanation.
- For fill_gap questions, include one clear blank using "____" in the stem; options should be possible gap completions and correct_answer must be the exact correct completion.
- For ordering questions, ask learners to sequence a process; options should be steps and correct_answer should state the correct ordered sequence.
- Distractors must be believable misconceptions from the topic, not obviously silly answers.
- Wrong options should reflect common mistakes, shallow interpretations, or over-simplifications from the supplied material.
- Avoid vague stems such as "Which is best?" unless the scenario gives enough context.
- Avoid absolutes like always/never unless they are central to the concept.
- For true_false questions, use exactly True and False as options.
- For single-answer question types, correct_answer must exactly match one of the options or one option letter.
- For multiple_choice only, correct_answer must contain multiple correct options separated by commas.
- Explanations should be concise, specific, and explain why the answer is correct by linking to the document concept and KSB alignment.
- Do not output separate distractor rationales because the API schema only stores one explanation field, but use strong plausible distractors.
- Do not mention "the source", "the PDF", "the slide", or "the document" in the learner-facing question text.
- Keep question text clear, but make the cognitive demand suitable for admin-approved assessment.
- Apply the editable embedded prompt from the admin UI when provided, but do not let it override safety, source-grounding, JSON schema, mixed question type, KSB alignment, or quality rules.
- Self-check before finalising: correct question count, single_choice has one best answer, multiple_choice has 2 or 3 correct answers, 4 options for MCQs, no duplicate stems, balanced answer letters where relevant, clear KSB-linked explanations.
"""

    try:
        from openai import OpenAI
    except ImportError:
        return JsonResponse({"error": "OpenAI Python package is not installed. Run pip install -r requirements.txt."}, status=503)

    try:
        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        response = client.responses.create(
            model=settings.OPENAI_MODEL,
            input=[
                {
                    "role": "system",
                    "content": "You generate assessment questions and return only structured JSON.",
                },
                {"role": "user", "content": prompt},
            ],
            text={
                "format": {
                    "type": "json_schema",
                    "name": "generated_quiz_questions",
                    "schema": schema,
                    "strict": True,
                }
            },
        )
        raw_text = response.output_text
        generated = json.loads(raw_text)
        questions = _normalise_generated_questions(generated, "single_choice")
        if not questions:
            return JsonResponse({"error": "AI did not return usable questions."}, status=502)
        return JsonResponse({
            "questions": questions,
            "source": {
                "model": settings.OPENAI_MODEL,
                "questionCount": len(questions),
                "difficulty": "mixed",
                "questionType": "mixed",
                "readableFiles": readable_files,
                "unreadableFiles": unreadable_files,
            },
        })
    except json.JSONDecodeError:
        return JsonResponse({"error": "AI returned invalid JSON."}, status=502)
    except Exception as exc:
        message = str(exc)
        if "insufficient_quota" in message or "exceeded your current quota" in message or "Error code: 429" in message:
            return JsonResponse({
                "error": "OpenAI quota is exhausted for the configured API key. Check billing/quota or replace OPENAI_API_KEY in backend .env.",
            }, status=429)
        if "invalid_api_key" in message or "Incorrect API key" in message:
            return JsonResponse({
                "error": "The configured OpenAI API key is invalid. Update OPENAI_API_KEY in backend .env.",
            }, status=401)
        return JsonResponse({"error": f"Question generation failed: {message}"}, status=502)


@csrf_exempt
@require_http_methods(["GET", "POST"])
def quizzes(request):
    _ensure_quiz_assessment_type_column()
    if request.method == "GET":
        query = request.GET.get("search", "").strip()
        status = request.GET.get("status", "all")
        assessment_type = request.GET.get("assessmentType", "").strip().lower()
        queryset = QuizPackage.objects.all()
        if assessment_type in ASSESSMENT_TYPES:
            queryset = queryset.filter(assessment_type=assessment_type)
        else:
            queryset = queryset.exclude(assessment_type="checkpoint")
        if status == "all":
            queryset = queryset.exclude(status="trash")
        else:
            if status == "pending":
                queryset = queryset.filter(status__in=["pending", "validating"])
            else:
                queryset = queryset.filter(status=status)
        if query:
            queryset = queryset.filter(title__icontains=query)
        return JsonResponse({"results": [_serialize_quiz(quiz) for quiz in queryset]})

    uploaded_file = request.FILES.get("file")
    if uploaded_file:
        extension = Path(uploaded_file.name).suffix.lower()
        package_type = _package_type_from_extension(extension)
        if package_type == "file":
            return JsonResponse({"error": "Unsupported quiz file type. Upload XML, SCORM ZIP, XLSX, XLSM or CSV."}, status=400)
        if package_type == "xml":
            valid, message, question_count, detected_title, parsed_questions, file_metadata = _validate_xml(uploaded_file)
        elif package_type == "scorm":
            valid, message, question_count, detected_title, parsed_questions, file_metadata = _validate_scorm(uploaded_file)
        else:
            valid, message, question_count, detected_title, parsed_questions, file_metadata = _extract_questions_from_upload(uploaded_file, package_type)

        title = request.POST.get("title") or detected_title or Path(uploaded_file.name).stem.replace("-", " ").title()
        programme = request.POST.get("programme") or file_metadata.get("programme", "")
        module = request.POST.get("module") or file_metadata.get("module", "")
        raw_programme_id = request.POST.get("programmeId")
        training_plan_id = raw_programme_id if _is_int_like(raw_programme_id) else _match_training_plan_id(programme, module, title)
        programme_id = _match_programme_catalogue_id(programme, module, title, raw_programme_id or training_plan_id)
        if training_plan_id and not programme:
            programme = _training_plan_programme_for_id(training_plan_id)
        week_value = request.POST.get("week") or request.POST.get("weekNumber") or title
        week_id = request.POST.get("weekId") or _build_week_id(programme_id, week_value)
        supplied_question_count = int(request.POST.get("questions") or 0)
        linked_courses = int(request.POST.get("linkedCourses") or (1 if programme_id else 0))
        final_question_count = len(parsed_questions) or question_count or supplied_question_count
        if package_type in {"excel", "csv"} and not parsed_questions:
            valid = False
            message = message or "File uploaded, but no questions were detected. Use columns like Question Title, Option 1-5, and Answer."
        default_question_type = _default_question_type_from_questions(parsed_questions, request.POST.get("questionType") or "single_choice")

        with transaction.atomic():
            quiz = QuizPackage.objects.create(
                title=title,
                programme_id=programme_id or "",
                module=module,
                programme=programme,
                version=request.POST.get("version", "v1.0"),
                questions=final_question_count,
                default_question_type=default_question_type,
                assessment_type=_clean_assessment_type(request.POST.get("assessmentType"), "quiz"),
                week_id=week_id,
                status=_clean_quiz_status(request.POST.get("status"), "draft"),
                package_type=package_type,
                uploaded_file=uploaded_file,
                file_name=uploaded_file.name,
                file_size=uploaded_file.size,
                schema_valid=valid,
                validation_message=message,
                mapped_components=int(request.POST.get("mappedComponents") or 0),
                author=request.POST.get("author", "Curriculum Team"),
                linked_courses=linked_courses,
            )
            _save_questions(quiz, parsed_questions)
        return JsonResponse(_serialize_quiz(quiz), status=201)

    try:
        payload = json.loads(request.body or "{}")
    except json.JSONDecodeError:
        return JsonResponse({"error": "Invalid JSON body"}, status=400)

    manual_title = payload.get("title", "Untitled Quiz")
    manual_programme = payload.get("programme", "")
    manual_module = payload.get("module", "")
    raw_manual_programme_id = payload.get("programmeId")
    manual_training_plan_id = raw_manual_programme_id if _is_int_like(raw_manual_programme_id) else _match_training_plan_id(manual_programme, manual_module, manual_title)
    manual_programme_id = _match_programme_catalogue_id(manual_programme, manual_module, manual_title, raw_manual_programme_id or manual_training_plan_id)
    if manual_training_plan_id and not manual_programme:
        manual_programme = _training_plan_programme_for_id(manual_training_plan_id)
    manual_week_value = payload.get("week") or payload.get("weekNumber") or manual_title
    manual_week_id = payload.get("weekId") or _build_week_id(manual_programme_id, manual_week_value)
    manual_question_type = _normalise_question_type(payload.get("questionType") or "single_choice")

    quiz = QuizPackage.objects.create(
        title=manual_title,
        programme_id=manual_programme_id or "",
        module=manual_module,
        programme=manual_programme,
        version=payload.get("version", "v1.0"),
        questions=int(payload.get("questions") or 0),
        default_question_type=manual_question_type,
        assessment_type=_clean_assessment_type(payload.get("assessmentType"), "quiz"),
        week_id=manual_week_id,
        status=_clean_quiz_status(payload.get("status"), "draft"),
        package_type=payload.get("packageType", "xml"),
        schema_valid=bool(payload.get("schemaValid", True)),
        mapped_components=int(payload.get("mappedComponents") or 0),
        author=payload.get("author", "Curriculum Team"),
        linked_courses=int(payload.get("linkedCourses") or (1 if manual_programme_id else 0)),
    )
    return JsonResponse(_serialize_quiz(quiz), status=201)


def _parse_grade_percent(value):
    if value is None:
        return None
    if isinstance(value, (int, float)):
        return round(float(value) * 100, 1) if 0 <= float(value) <= 1 else round(float(value), 1)
    match = re.search(r"-?\d+(?:\.\d+)?", str(value))
    return float(match.group(0)) if match else None


def _format_grade(value):
    percent = _parse_grade_percent(value)
    if percent is None:
        return ""
    return f"{int(percent) if percent == int(percent) else percent}%"


def _attempt_sort_key(attempt):
    submitted_at = str(attempt.get("submittedAt") or "")
    try:
        attempt_number = int(attempt.get("attempt") or 0)
    except (TypeError, ValueError):
        attempt_number = 0
    return submitted_at, attempt_number


def _answer_label(answer_id, answer_lookup):
    if answer_id is None:
        return None
    if isinstance(answer_id, list):
        labels = [_answer_label(item, answer_lookup) for item in answer_id]
        return ", ".join(label for label in labels if label) or None
    try:
        key = int(answer_id)
    except (TypeError, ValueError):
        return str(answer_id)
    return answer_lookup.get(key) or str(answer_id)


def _quiz_question_lookup(quiz_id):
    questions = QuizQuestion.objects.filter(quiz_id=quiz_id).prefetch_related("answers")
    question_lookup = {}
    answer_lookup = {}
    for question in questions:
        question_lookup[question.id] = {
            "text": question.question_text,
            "type": question.question_type,
            "possible": question.points,
        }
        for answer in question.answers.all():
            answer_lookup[answer.id] = answer.answer_text
    return question_lookup, answer_lookup


def _serialize_attempt_for_student(attempt, question_lookup, answer_lookup):
    questions = []
    for index, question in enumerate(attempt.get("questions") or [], start=1):
        question_id = question.get("questionId") or question.get("id")
        try:
            question_id = int(question_id)
        except (TypeError, ValueError):
            pass
        question_meta = question_lookup.get(question_id) or {}
        chosen_answer = question.get("chosenAnswer")
        correct_answer = question.get("correctAnswer")
        if chosen_answer is None:
            chosen_answer = _answer_label(question.get("chosenAnswerId"), answer_lookup)
        if correct_answer is None:
            correct_answer = _answer_label(question.get("correctAnswerId"), answer_lookup)
        questions.append({
            "number": index,
            "text": question.get("text") or question_meta.get("text") or "",
            "type": question.get("type") or question_meta.get("type") or "",
            "chosenAnswer": chosen_answer,
            "correctAnswer": correct_answer,
            "correct": bool(question.get("correct")),
            "earned": question.get("earned"),
            "possible": question.get("possible") or question_meta.get("possible"),
        })

    grade_percent = _parse_grade_percent(attempt.get("grade"))
    achieved_score = attempt.get("achievedScore")
    total_score = attempt.get("totalScore")
    score = attempt.get("Score")
    if not score and achieved_score is not None and total_score is not None:
        score = f"{achieved_score}/{total_score}"

    return {
        "attempt": attempt.get("attempt"),
        "grade": attempt.get("grade") if isinstance(attempt.get("grade"), str) and "%" in attempt.get("grade") else _format_grade(attempt.get("grade")),
        "gradePercent": grade_percent,
        "score": score or "",
        "passed": bool(attempt.get("passed")),
        "submittedAt": attempt.get("submittedAt") or "",
        "startedAt": attempt.get("startedAt") or "",
        "timeTaken": attempt.get("timeTaken") or "",
        "reportedTime": attempt.get("reportedTime") or "",
        "week": attempt.get("week") or "",
        "module": attempt.get("module") or "",
        "feedback": attempt.get("feedback") or "",
        "ksbs": attempt.get("ksbs") if isinstance(attempt.get("ksbs"), list) else [],
        "questions": questions,
    }


@require_http_methods(["GET"])
def quiz_students(request, pk):
    try:
        quiz = QuizPackage.objects.get(pk=pk)
    except QuizPackage.DoesNotExist:
        raise Http404("Quiz not found")

    try:
        question_lookup, answer_lookup = _quiz_question_lookup(pk)
        active_users = ActiveUser.objects.all().only(
            "id",
            "username",
            "email",
            "programme",
            "cohort",
            "group",
            "training_plan_progress",
        )
        students = []
        all_grade_percents = []
        for learner in active_users:
            history = learner.training_plan_progress if isinstance(learner.training_plan_progress, list) else []
            quiz_attempts = [
                attempt for attempt in history
                if attempt.get("kind") == "quiz" and str(attempt.get("quizId")) == str(pk)
            ]
            if not quiz_attempts:
                continue

            serialized_attempts = [
                _serialize_attempt_for_student(attempt, question_lookup, answer_lookup)
                for attempt in sorted(quiz_attempts, key=_attempt_sort_key, reverse=True)
            ]
            latest_attempt = serialized_attempts[0]
            best_grade = max(
                (attempt["gradePercent"] for attempt in serialized_attempts if attempt["gradePercent"] is not None),
                default=None,
            )
            if best_grade is not None:
                all_grade_percents.append(best_grade)

            students.append({
                "id": learner.id,
                "name": learner.username or learner.email or f"Learner {learner.id}",
                "email": learner.email or "",
                "programme": learner.programme or "",
                "cohort": learner.cohort or "",
                "group": learner.group or "",
                "attemptCount": len(serialized_attempts),
                "bestGrade": best_grade,
                "latestAttempt": latest_attempt,
                "attempts": serialized_attempts,
            })
    except DatabaseError as exc:
        return JsonResponse({"error": f"Database error: {exc}"}, status=502)

    students.sort(key=lambda item: (item["latestAttempt"].get("submittedAt") or "", item["name"]), reverse=True)
    passed_count = sum(1 for student in students if student["latestAttempt"].get("passed"))
    average_best = round(sum(all_grade_percents) / len(all_grade_percents), 1) if all_grade_percents else None

    return JsonResponse({
        "quiz": _serialize_quiz(quiz),
        "summary": {
            "students": len(students),
            "attempts": sum(student["attemptCount"] for student in students),
            "passed": passed_count,
            "averageBest": average_best,
        },
        "students": students,
    })


@csrf_exempt
@require_http_methods(["PATCH", "DELETE"])
def quiz_detail(request, pk):
    _ensure_quiz_assessment_type_column()
    try:
        quiz = QuizPackage.objects.get(pk=pk)
    except QuizPackage.DoesNotExist:
        raise Http404

    if request.method == "DELETE":
        quiz.delete()
        return JsonResponse({"deleted": True})

    try:
        payload = json.loads(request.body or "{}")
    except json.JSONDecodeError:
        return JsonResponse({"error": "Invalid JSON body"}, status=400)

    if "status" in payload:
        status = _clean_quiz_status(payload["status"], "")
        if not status:
            return JsonResponse({"error": "Unsupported quiz status"}, status=400)
        quiz.status = status

    for field in ["title", "module", "programme", "version", "author"]:
        if field in payload:
            setattr(quiz, field, payload[field])
    for field, attr in [
        ("shortDescription", "short_description"),
        ("lessonContent", "lesson_content"),
        ("timeUnit", "time_unit"),
        ("quizStyle", "quiz_style"),
    ]:
        if field in payload:
            setattr(quiz, attr, payload[field] or "")
    if "questionType" in payload:
        quiz.default_question_type = payload["questionType"] or "single_choice"
    if "assessmentType" in payload:
        quiz.assessment_type = _clean_assessment_type(payload["assessmentType"], quiz.assessment_type or "quiz")
    if "programmeId" in payload:
        quiz.programme_id = _match_programme_catalogue_id(quiz.programme, quiz.module, quiz.title, payload["programmeId"]) or ""
    if "weekId" in payload:
        quiz.week_id = payload["weekId"] or ""
    elif any(field in payload for field in ["week", "weekNumber", "title", "programmeId", "programme", "module"]):
        quiz.week_id = _build_week_id(quiz.programme_id, payload.get("week") or payload.get("weekNumber") or quiz.title)
    for field, attr in [
        ("questions", "questions"),
        ("mappedComponents", "mapped_components"),
        ("linkedCourses", "linked_courses"),
        ("duration", "duration"),
        ("passingGrade", "passing_grade"),
        ("retakePointsCut", "retake_points_cut"),
    ]:
        if field in payload:
            setattr(quiz, attr, int(payload[field] or 0))
    for field, attr in [
        ("randomizeQuestions", "randomize_questions"),
        ("randomizeAnswers", "randomize_answers"),
        ("showCorrectAnswer", "show_correct_answer"),
        ("attemptHistory", "attempt_history"),
        ("retakeAfterPass", "retake_after_pass"),
        ("limitAttempts", "limit_attempts"),
    ]:
        if field in payload:
            setattr(quiz, attr, bool(payload[field]))
    if "schemaValid" in payload:
        quiz.schema_valid = bool(payload["schemaValid"])
    if quiz.status == "published" and not quiz.published_at:
        quiz.published_at = timezone.now()
    quiz.save()
    return JsonResponse(_serialize_quiz(quiz))


@csrf_exempt
@require_http_methods(["GET", "PATCH"])
def quiz_course_links(request, pk):
    _ensure_quiz_assessment_type_column()
    try:
        quiz = QuizPackage.objects.get(pk=pk)
    except QuizPackage.DoesNotExist:
        raise Http404

    try:
        _ensure_quiz_course_links_table()
    except Exception as exc:
        return JsonResponse({"error": f"Could not prepare course links table: {exc}"}, status=500)

    programme = quiz.programme or _training_plan_programme_for_id(quiz.programme_id)
    if request.method == "PATCH":
        try:
            payload = json.loads(request.body or "{}")
        except json.JSONDecodeError:
            return JsonResponse({"error": "Invalid JSON body"}, status=400)

        requested_ids = payload.get("trainingPlanIds", [])
        if not isinstance(requested_ids, list):
            return JsonResponse({"error": "trainingPlanIds must be a list"}, status=400)

        valid_courses = _training_plan_courses_for_programme(programme)
        valid_ids = {course["id"] for course in valid_courses}
        selected_ids = []
        for raw_id in requested_ids:
            try:
                plan_id = int(raw_id)
            except (TypeError, ValueError):
                continue
            if plan_id in valid_ids and plan_id not in selected_ids:
                selected_ids.append(plan_id)

        with transaction.atomic():
            with connection.cursor() as cursor:
                cursor.execute("delete from curriculum.quiz_course_links where quiz_id = %s", [quiz.id])
                for plan_id in selected_ids:
                    cursor.execute(
                        """
                        insert into curriculum.quiz_course_links (quiz_id, training_plan_id)
                        values (%s, %s)
                        on conflict (quiz_id, training_plan_id) do nothing
                        """,
                        [quiz.id, plan_id],
                    )
            quiz.linked_courses = len(selected_ids)
            if not quiz.programme_id:
                quiz.programme_id = _match_programme_catalogue_id(quiz.programme, quiz.module, quiz.title)
                quiz.save(update_fields=["linked_courses", "programme_id", "updated_at"])
            else:
                quiz.save(update_fields=["linked_courses", "updated_at"])

    courses = _training_plan_courses_for_programme(programme)
    selected_ids = set(_quiz_course_link_ids(quiz.id))
    if not selected_ids and _is_int_like(quiz.programme_id):
        selected_ids = {int(quiz.programme_id)}

    return JsonResponse({
        "programme": programme or "",
        "selectedIds": sorted(selected_ids),
        "courses": [
            {**course, "selected": course["id"] in selected_ids}
            for course in courses
        ],
        "quiz": _serialize_quiz(quiz),
    })


def _safe_download_name(quiz, extension):
    base_name = Path(quiz.file_name or "").stem or quiz.title or f"quiz-{quiz.id}"
    safe_name = re.sub(r"[^A-Za-z0-9._-]+", "-", base_name).strip(".-") or f"quiz-{quiz.id}"
    return f"{safe_name}.{extension}"


def _export_questions(quiz):
    return (
        QuizQuestion.objects
        .filter(quiz_id=quiz.id, is_archived=False)
        .prefetch_related("answers")
        .order_by("sort_order", "id")
    )


def _quiz_csv_response(quiz):
    output = io.StringIO()
    writer = csv.writer(output)
    writer.writerow([
        "question_number",
        "question_type",
        "question",
        "option_a",
        "option_b",
        "option_c",
        "option_d",
        "option_e",
        "correct_answer",
        "feedback",
        "programme",
        "module",
    ])

    for index, question in enumerate(_export_questions(quiz), start=1):
        answers = list(question.answers.all())
        option_texts = [answer.answer_text for answer in answers[:5]]
        option_texts += [""] * (5 - len(option_texts))
        correct_letters = [
            chr(65 + answer_index)
            for answer_index, answer in enumerate(answers[:5])
            if answer.is_correct
        ]
        writer.writerow([
            index,
            question.question_type,
            question.question_text,
            *option_texts[:5],
            ",".join(correct_letters),
            question.explanation,
            quiz.programme,
            quiz.module,
        ])

    response = HttpResponse(output.getvalue(), content_type="text/csv; charset=utf-8")
    response["Content-Disposition"] = f'attachment; filename="{_safe_download_name(quiz, "csv")}"'
    return response


def _quiz_xml_response(quiz):
    root = ElementTree.Element("quiz", {
        "id": str(quiz.id),
        "version": quiz.version or "v1.0",
        "status": quiz.status,
    })
    metadata = ElementTree.SubElement(root, "metadata")
    for tag, value in [
        ("title", quiz.title),
        ("module", quiz.module),
        ("programme", quiz.programme),
        ("package_type", quiz.package_type),
    ]:
        ElementTree.SubElement(metadata, tag).text = value or ""

    questions_node = ElementTree.SubElement(root, "questions", {
        "count": str(QuizQuestion.objects.filter(quiz_id=quiz.id, is_archived=False).count())
    })
    for index, question in enumerate(_export_questions(quiz), start=1):
        question_node = ElementTree.SubElement(questions_node, "question", {
            "id": str(question.id),
            "number": str(index),
            "type": question.question_type,
            "points": str(question.points),
        })
        ElementTree.SubElement(question_node, "text").text = question.question_text
        answers_node = ElementTree.SubElement(question_node, "answers")
        for answer_index, answer in enumerate(question.answers.all(), start=1):
            answer_node = ElementTree.SubElement(answers_node, "answer", {
                "id": str(answer.id),
                "letter": chr(64 + answer_index),
                "correct": "true" if answer.is_correct else "false",
            })
            answer_node.text = answer.answer_text
        ElementTree.SubElement(question_node, "feedback").text = question.explanation or ""

    content = ElementTree.tostring(root, encoding="utf-8", xml_declaration=True)
    response = HttpResponse(content, content_type="application/xml; charset=utf-8")
    response["Content-Disposition"] = f'attachment; filename="{_safe_download_name(quiz, "xml")}"'
    return response


def _scorm_uploaded_file_readable(quiz):
    """True only when the quiz has an uploaded package that actually exists on disk."""
    if not quiz.uploaded_file:
        return False
    try:
        return quiz.uploaded_file.storage.exists(quiz.uploaded_file.name)
    except Exception:
        return False


def _scorm_index_html(quiz):
    questions = [
        {
            "text": question.question_text,
            "type": question.question_type,
            "feedback": question.explanation,
            "answers": [
                {"text": answer.answer_text, "isCorrect": answer.is_correct}
                for answer in question.answers.all()
            ],
        }
        for question in _export_questions(quiz)
    ]
    quiz_payload = json.dumps({
        "title": quiz.title,
        "programme": quiz.programme,
        "module": quiz.module,
        "questions": questions,
    })
    return f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{quiz.title}</title>
  <style>
    body {{ margin: 0; font-family: Arial, sans-serif; background: #f8fafc; color: #0f172a; }}
    main {{ max-width: 920px; margin: 0 auto; padding: 32px 20px; }}
    .card {{ background: #fff; border: 1px solid #dbe3ee; border-radius: 18px; padding: 24px; margin-bottom: 18px; }}
    .answer {{ border: 1px solid #dbe3ee; border-radius: 12px; padding: 12px 14px; margin-top: 10px; }}
    .correct {{ border-color: #34d399; background: #ecfdf5; color: #047857; }}
    .feedback {{ margin-top: 14px; border: 1px solid #ddd2ff; background: #fbf9ff; border-radius: 12px; padding: 12px; color: #3f2f73; }}
  </style>
</head>
<body>
  <main id="app"></main>
  <script>
    const quiz = {quiz_payload};
    const app = document.getElementById('app');
    app.innerHTML = `<h1>${{quiz.title}}</h1><p>${{quiz.questions.length}} questions - ${{quiz.programme || ''}}</p>` +
      quiz.questions.map((question, index) => `
        <section class="card">
          <h2>${{index + 1}}. ${{question.text}}</h2>
          ${{question.answers.map((answer, answerIndex) => `
            <div class="answer ${{answer.isCorrect ? 'correct' : ''}}">
              <strong>${{String.fromCharCode(65 + answerIndex)}}.</strong> ${{answer.text}}
              ${{answer.isCorrect ? '<strong style="float:right">Correct</strong>' : ''}}
            </div>
          `).join('')}}
          ${{question.feedback ? `<div class="feedback"><strong>Feedback</strong><br>${{question.feedback}}</div>` : ''}}
        </section>
      `).join('');
  </script>
</body>
</html>
"""


def _quiz_scorm_response(quiz):
    if _scorm_uploaded_file_readable(quiz):
        filename = quiz.file_name or _safe_download_name(quiz, "zip")
        if not filename.lower().endswith(".zip"):
            filename = _safe_download_name(quiz, "zip")
        return FileResponse(quiz.uploaded_file.open("rb"), as_attachment=True, filename=filename)

    index_html = _scorm_index_html(quiz)
    manifest = f"""<?xml version="1.0" encoding="UTF-8"?>
<manifest identifier="quiz-{quiz.id}" version="1.2"
  xmlns="http://www.imsproject.org/xsd/imscp_rootv1p1p2"
  xmlns:adlcp="http://www.adlnet.org/xsd/adlcp_rootv1p2">
  <metadata>
    <schema>ADL SCORM</schema>
    <schemaversion>1.2</schemaversion>
  </metadata>
  <organizations default="org-{quiz.id}">
    <organization identifier="org-{quiz.id}">
      <title>{quiz.title}</title>
      <item identifier="item-{quiz.id}" identifierref="res-{quiz.id}">
        <title>{quiz.title}</title>
      </item>
    </organization>
  </organizations>
  <resources>
    <resource identifier="res-{quiz.id}" type="webcontent" adlcp:scormtype="sco" href="index.html">
      <file href="index.html" />
    </resource>
  </resources>
</manifest>
"""
    archive = io.BytesIO()
    with zipfile.ZipFile(archive, "w", compression=zipfile.ZIP_DEFLATED) as package:
        package.writestr("imsmanifest.xml", manifest)
        package.writestr("index.html", index_html)
    archive.seek(0)

    response = HttpResponse(archive.getvalue(), content_type="application/zip")
    response["Content-Disposition"] = f'attachment; filename="{_safe_download_name(quiz, "zip")}"'
    return response


@require_http_methods(["GET"])
def quiz_download(request, pk):
    _ensure_quiz_assessment_type_column()
    try:
        quiz = QuizPackage.objects.get(pk=pk)
    except QuizPackage.DoesNotExist:
        raise Http404

    package_type = (quiz.package_type or "xml").lower()
    if package_type == "scorm":
        return _quiz_scorm_response(quiz)
    if package_type in {"csv", "excel"}:
        return _quiz_csv_response(quiz)
    return _quiz_xml_response(quiz)


@require_http_methods(["GET"])
@xframe_options_sameorigin
def quiz_scorm_launch(request, pk, asset_path=""):
    _ensure_quiz_assessment_type_column()
    try:
        quiz = QuizPackage.objects.get(pk=pk)
    except QuizPackage.DoesNotExist:
        raise Http404

    if quiz.package_type != "scorm":
        raise Http404

    # No real package on disk (or it never uploaded): serve a generated preview
    # built from the quiz's saved questions instead of returning a 500/404 that
    # the browser renders as an empty, broken iframe.
    if not _scorm_uploaded_file_readable(quiz):
        if asset_path and normpath(asset_path.replace("\\", "/")).lstrip("/") not in ("", ".", "index.html"):
            raise Http404
        return HttpResponse(_scorm_index_html(quiz), content_type="text/html; charset=utf-8")

    destination, launch_path = _ensure_scorm_extracted(quiz)
    requested_path = asset_path or launch_path
    target = _safe_scorm_member_path(destination, requested_path)
    if target is None or not target.is_file():
        raise Http404

    content_type = mimetypes.guess_type(target.name)[0] or "application/octet-stream"
    return FileResponse(target.open("rb"), content_type=content_type)


def _quiz_preview_payload(quiz, include_archived=False):
    questions = []
    queryset = QuizQuestion.objects.filter(quiz_id=quiz.id)
    if not include_archived:
        queryset = queryset.filter(is_archived=False)
    for question in queryset.prefetch_related("answers"):
        questions.append({
            "id": question.id,
            "text": question.question_text,
            "questionType": question.question_type,
            "explanation": question.explanation,
            "isArchived": question.is_archived,
            "answers": [
                {
                    "id": answer.id,
                    "text": answer.answer_text,
                    "isCorrect": answer.is_correct,
                }
                for answer in question.answers.all()
            ],
        })

    return {
        "quiz": _serialize_quiz(quiz),
        "questions": questions,
    }


@require_http_methods(["GET"])
def quiz_preview(request, pk):
    _ensure_quiz_assessment_type_column()
    try:
        quiz = QuizPackage.objects.get(pk=pk)
    except QuizPackage.DoesNotExist:
        raise Http404

    return JsonResponse(_quiz_preview_payload(quiz))


@require_http_methods(["GET"])
def question_bank(request):
    _ensure_quiz_assessment_type_column()
    selected_programme = request.GET.get("programme", "all").strip()
    query = request.GET.get("search", "").strip()
    base_queryset = (
        QuizQuestion.objects
        .filter(is_archived=False)
        .exclude(quiz__status="trash")
        .select_related("quiz")
        .prefetch_related("answers")
    )
    if query:
        base_queryset = base_queryset.filter(question_text__icontains=query)

    programme_stats = {}
    for programme in _training_plan_programmes():
        programme_stats[programme["name"]] = {
            "name": programme["name"],
            "questionCount": 0,
            "quizCount": 0,
            "trainingPlanRows": programme["trainingPlanRows"],
        }

    all_questions = list(
        base_queryset.order_by("quiz__programme_id", "quiz__module", "quiz__title", "sort_order", "id")
    )
    available_quizzes = list(QuizPackage.objects.exclude(status="trash").order_by("programme_id", "title"))
    plan_map = _training_plan_programme_map({
        *{question.quiz.programme_id for question in all_questions},
        *{quiz.programme_id for quiz in available_quizzes},
    })
    matched_plan_cache = {}

    def resolve_programme(quiz):
        if quiz.programme_id and plan_map.get(quiz.programme_id):
            return plan_map[quiz.programme_id], plan_map[quiz.programme_id]
        cache_key = quiz.id
        if cache_key not in matched_plan_cache:
            matched_plan_id = _match_training_plan_id(quiz.programme, quiz.module, quiz.title)
            if matched_plan_id and matched_plan_id not in plan_map:
                plan_map.update(_training_plan_programme_map([matched_plan_id]))
            matched_plan_cache[cache_key] = plan_map.get(matched_plan_id, "") if matched_plan_id else ""
        if matched_plan_cache[cache_key]:
            return matched_plan_cache[cache_key], matched_plan_cache[cache_key]
        return "__unassigned__", "Unassigned"

    quiz_ids_by_programme = {}
    question_programme_info = {}
    question_usage_by_signature = {}
    for question in all_questions:
        key, programme_name = resolve_programme(question.quiz)
        question_programme_info[question.id] = (key, programme_name)
        if key not in programme_stats:
            programme_stats[key] = {
                "name": programme_name,
                "questionCount": 0,
                "quizCount": 0,
                "trainingPlanRows": 0,
            }
        programme_stats[key]["questionCount"] += 1
        quiz_ids_by_programme.setdefault(key, set()).add(question.quiz_id)
        signature = (
            key,
            question.question_type,
            " ".join((question.question_text or "").split()).strip().lower(),
        )
        linked_quizzes = question_usage_by_signature.setdefault(signature, {})
        linked_quizzes[question.quiz_id] = {
            "id": question.quiz_id,
            "title": question.quiz.title,
            "module": question.quiz.module,
            "status": question.quiz.status,
        }

    questions = []
    for question in all_questions:
        key, programme_name = question_programme_info[question.id]

        if selected_programme != "all" and selected_programme != key:
            continue

        signature = (
            key,
            question.question_type,
            " ".join((question.question_text or "").split()).strip().lower(),
        )

        questions.append({
            "id": question.id,
            "text": question.question_text,
            "questionType": question.question_type,
            "explanation": question.explanation,
            "programme": programme_name,
            "programmeKey": key,
            "module": question.quiz.module,
            "quizId": question.quiz_id,
            "quizTitle": question.quiz.title,
            "quizStatus": question.quiz.status,
            "linkedQuizzes": list(question_usage_by_signature.get(signature, {}).values()),
            "answers": [
                {
                    "id": answer.id,
                    "text": answer.answer_text,
                    "isCorrect": answer.is_correct,
                }
                for answer in question.answers.all()
            ],
        })

    for key, quiz_ids in quiz_ids_by_programme.items():
        programme_stats[key]["quizCount"] = len(quiz_ids)

    quizzes = []
    for quiz in available_quizzes:
        key, programme_name = resolve_programme(quiz)
        quizzes.append({
            "id": quiz.id,
            "title": quiz.title,
            "module": quiz.module,
            "programme": programme_name,
            "programmeKey": key,
            "questions": quiz.questions,
            "status": quiz.status,
        })

    programmes = sorted(
        programme_stats.values(),
        key=lambda item: (item["name"] == "Unassigned", item["name"].lower()),
    )
    return JsonResponse({
        "programmes": programmes,
        "questions": questions,
        "quizzes": quizzes,
        "totalQuestions": sum(programme["questionCount"] for programme in programmes),
    })


@csrf_exempt
@require_http_methods(["POST"])
def add_question_to_quiz(request, question_id):
    _ensure_quiz_assessment_type_column()
    try:
        payload = json.loads(request.body or "{}")
    except json.JSONDecodeError:
        return JsonResponse({"error": "Invalid JSON body"}, status=400)

    target_quiz_id = payload.get("quizId")
    if not target_quiz_id:
        return JsonResponse({"error": "quizId is required"}, status=400)

    try:
        source_question = QuizQuestion.objects.select_related("quiz").prefetch_related("answers").get(pk=question_id)
    except QuizQuestion.DoesNotExist:
        raise Http404

    with transaction.atomic():
        try:
            target_quiz = QuizPackage.objects.select_for_update().get(pk=target_quiz_id)
        except QuizPackage.DoesNotExist:
            return JsonResponse({"error": "Target quiz not found"}, status=404)

        next_sort_order = (
            QuizQuestion.objects
            .filter(quiz=target_quiz)
            .aggregate(max_order=Max("sort_order"))
            .get("max_order")
        )
        next_sort_order = 0 if next_sort_order is None else next_sort_order + 1

        new_question = QuizQuestion.objects.create(
            quiz=target_quiz,
            question_text=source_question.question_text,
            question_type=source_question.question_type,
            points=source_question.points,
            sort_order=next_sort_order,
            explanation=source_question.explanation,
            is_archived=False,
        )
        for answer in source_question.answers.all():
            QuizAnswer.objects.create(
                question=new_question,
                answer_text=answer.answer_text,
                is_correct=answer.is_correct,
                sort_order=answer.sort_order,
            )

        target_quiz.questions = QuizQuestion.objects.filter(quiz=target_quiz, is_archived=False).count()
        target_quiz.save(update_fields=["questions", "updated_at"])

    target_quiz.refresh_from_db()
    return JsonResponse({
        "question": {
            "id": new_question.id,
            "text": new_question.question_text,
            "quizId": target_quiz.id,
            "quizTitle": target_quiz.title,
        },
        "quiz": _serialize_quiz(target_quiz),
    }, status=201)


@csrf_exempt
def quiz_questions(request, pk):
    _ensure_quiz_assessment_type_column()
    if request.method not in {"GET", "POST", "PATCH"}:
        return JsonResponse({"error": "Method not allowed"}, status=405)

    try:
        quiz = QuizPackage.objects.get(pk=pk)
    except QuizPackage.DoesNotExist:
        raise Http404

    if request.method == "GET":
        return JsonResponse(_quiz_preview_payload(quiz, include_archived=True))

    try:
        payload = json.loads(request.body or "{}")
    except json.JSONDecodeError:
        return JsonResponse({"error": "Invalid JSON body"}, status=400)

    questions = payload.get("questions", [])
    if not isinstance(questions, list):
        return JsonResponse({"error": "questions must be a list"}, status=400)

    existing_questions = {
        question.id: question
        for question in QuizQuestion.objects.filter(quiz_id=quiz.id).prefetch_related("answers")
    }
    remove_missing = bool(payload.get("removeMissing"))
    if remove_missing and len(questions) == 1 and len(existing_questions) > 1:
        return JsonResponse({
            "error": "Refusing to replace a multi-question quiz with one question. Reload the editor and try again.",
        }, status=409)

    seen_question_ids = set()
    with transaction.atomic():
        for question_index, question in enumerate(questions):
            question_id = question.get("id")
            question_record = existing_questions.get(question_id)
            if question_record:
                question_record.question_text = question.get("text", "").strip() or "Untitled question"
                question_record.question_type = question.get("questionType") or quiz.default_question_type
                question_record.explanation = question.get("explanation", "")
                question_record.is_archived = bool(question.get("isArchived"))
                question_record.sort_order = question_index
                question_record.save(update_fields=["question_text", "question_type", "explanation", "is_archived", "sort_order"])
                question_record.answers.all().delete()
                seen_question_ids.add(question_record.id)
            else:
                question_record = QuizQuestion.objects.create(
                    quiz=quiz,
                    question_text=question.get("text", "").strip() or "Untitled question",
                    question_type=question.get("questionType") or quiz.default_question_type,
                    explanation=question.get("explanation", ""),
                    is_archived=bool(question.get("isArchived")),
                    sort_order=question_index,
                )
                seen_question_ids.add(question_record.id)

            for answer_index, answer in enumerate(question.get("answers", [])):
                answer_text = answer.get("text", "").strip()
                if not answer_text:
                    continue
                QuizAnswer.objects.create(
                    question=question_record,
                    answer_text=answer_text,
                    is_correct=bool(answer.get("isCorrect")),
                    sort_order=answer_index,
                )

        if remove_missing:
            QuizQuestion.objects.filter(quiz_id=quiz.id).exclude(id__in=seen_question_ids).delete()
        else:
            next_order = len(questions)
            for existing_id, existing_question in existing_questions.items():
                if existing_id in seen_question_ids:
                    continue
                existing_question.sort_order = next_order
                existing_question.save(update_fields=["sort_order"])
                next_order += 1

        quiz.questions = QuizQuestion.objects.filter(quiz_id=quiz.id, is_archived=False).count()
        quiz.save(update_fields=["questions", "updated_at"])

    quiz.refresh_from_db()
    return JsonResponse(_quiz_preview_payload(quiz))
